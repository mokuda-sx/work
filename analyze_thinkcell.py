"""
Analyze think-cell .potx template files and extract slide/shape metadata.
Uses XML-level access to handle think-cell custom shapes that break python-pptx shape factory.
Patches content type to open .potx as .pptx.
"""
import json
import io
import sys
import zipfile
from pathlib import Path
from pptx import Presentation
from lxml import etree

BASE_DIR = Path(r"C:\ProgramData\test\work\refs\thinkcell")
OUTPUT_FILE = BASE_DIR / "_analysis_summary.json"

PRIORITY_FILES = [
    "processes, flow charts, phases/Process, Flow Chart, Phase.potx",
    "diagrams/funnels/Funnel.potx",
    "matrices, swot analyses/Matrix, SWOT Analysis.potx",
    "dashboards, statistics/Dashboard, Statistic.potx",
    "infographics, numbers/Infographic, Number.potx",
    "tables/Table.potx",
    "timelines, milestones, project planning/Timeline, Milestone, Project Planning.potx",
    "mental models, frameworks, concepts/Mental Model, Framework, Concept.potx",
]

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

SHAPE_TAGS = ("}sp", "}pic", "}grpSp", "}graphicFrame", "}cxnSp")


def emu_to_inches(emu_val):
    if emu_val is None:
        return None
    try:
        return round(int(emu_val) / 914400, 2)
    except (ValueError, TypeError):
        return None


def open_potx(filepath):
    """Open a .potx file by patching content type to presentation."""
    buf = io.BytesIO()
    with zipfile.ZipFile(str(filepath), "r") as zin:
        with zipfile.ZipFile(buf, "w") as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    data = data.replace(
                        b"presentationml.template.main+xml",
                        b"presentationml.presentation.main+xml",
                    )
                zout.writestr(item, data)
    buf.seek(0)
    return Presentation(buf)


def get_cNvPr(shape_elm):
    """Find cNvPr element for any shape type."""
    for child in shape_elm:
        ctag = child.tag.split("}")[-1]
        if ctag.startswith("nv"):
            for cc in child:
                if cc.tag.split("}")[-1] == "cNvPr":
                    return cc
    return None


def get_shape_text(shape_elm):
    """Extract text from shape's txBody."""
    texts = []
    for t_el in shape_elm.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}t"):
        if t_el.text:
            texts.append(t_el.text)
    return " ".join(texts).strip()


def get_shape_position(shape_elm):
    """Extract position/size from spPr/xfrm."""
    xfrm = shape_elm.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm")
    if xfrm is None:
        return {}
    off = xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}off")
    ext = xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ext")
    result = {}
    if off is not None:
        result["left"] = emu_to_inches(off.get("x"))
        result["top"] = emu_to_inches(off.get("y"))
    if ext is not None:
        result["width"] = emu_to_inches(ext.get("cx"))
        result["height"] = emu_to_inches(ext.get("cy"))
    return result


def analyze_potx(filepath):
    rel_path = str(filepath.relative_to(BASE_DIR)).replace(chr(92), "/")
    category = rel_path.split("/")[0]

    try:
        prs = open_potx(filepath)
    except Exception as e:
        return {"category": category, "file": rel_path, "error": str(e),
                "slide_count": 0, "slides": []}

    sw = emu_to_inches(prs.slide_width)
    sh = emu_to_inches(prs.slide_height)

    slides_info = []
    for idx, slide in enumerate(prs.slides):
        # Layout name via python-pptx (safe)
        layout_name = ""
        try:
            layout_name = slide.slide_layout.name
        except Exception:
            pass

        # Use XML for shapes
        slide_xml = slide._element
        sp_tree = slide_xml.find(".//p:cSld/p:spTree", NSMAP)
        if sp_tree is None:
            slides_info.append({"index": idx, "layout_name": layout_name,
                               "title": "", "shape_count": 0, "shape_names": [], "shapes": []})
            continue

        all_children = list(sp_tree)
        shape_elms = [s for s in all_children if any(s.tag.endswith(t) for t in SHAPE_TAGS)]

        title_text = ""
        shape_names = []
        shape_details = []

        for shp_el in shape_elms:
            tag = shp_el.tag.split("}")[-1]
            cnvpr = get_cNvPr(shp_el)
            name = cnvpr.get("name", "") if cnvpr is not None else ""

            pos = get_shape_position(shp_el)
            text = get_shape_text(shp_el)

            # Detect title
            if not title_text:
                # Check if this is a title placeholder
                ph = shp_el.find(".//p:nvSpPr/p:nvPr/p:ph", NSMAP)
                if ph is not None and ph.get("type", "") in ("title", "ctrTitle"):
                    title_text = text
                elif "title" in name.lower() and text:
                    title_text = text

            detail = {"name": name, "type": tag}
            detail.update(pos)
            if text:
                if len(text) > 80:
                    text = text[:77] + "..."
                detail["text"] = text

            shape_details.append(detail)
            shape_names.append("{} ({})".format(name, tag))

        slides_info.append({
            "index": idx,
            "layout_name": layout_name,
            "title": title_text[:100] if title_text else "",
            "shape_count": len(shape_elms),
            "shape_names": shape_names,
            "shapes": shape_details,
        })

    return {
        "category": category,
        "file": rel_path,
        "slide_size": "{} x {} inches".format(sw, sh),
        "slide_count": len(prs.slides),
        "slides": slides_info,
    }


def main():
    all_results = []
    all_potx = sorted(BASE_DIR.rglob("*.potx"))

    priority_paths = []
    for rel in PRIORITY_FILES:
        full = BASE_DIR / rel
        if full.exists():
            priority_paths.append(full)
        else:
            sys.stderr.write("[WARN] Not found: {}\n".format(rel))

    remaining = [p for p in all_potx if p not in priority_paths]
    ordered = priority_paths + remaining

    print("Found {} .potx files total, {} priority files.".format(len(all_potx), len(priority_paths)))
    print("=" * 70)

    for filepath in ordered:
        rel = str(filepath.relative_to(BASE_DIR)).replace(chr(92), "/")
        print("\nAnalyzing: {}".format(rel))
        result = analyze_potx(filepath)
        all_results.append(result)
        print(json.dumps(result, indent=2, ensure_ascii=True))
        print("---")

    with open(str(OUTPUT_FILE), "w", encoding="utf-8") as f:
        json.dump(all_results, f, indent=2, ensure_ascii=False)

    print("\nWrote combined summary to: {}".format(OUTPUT_FILE))
    print("Total files analyzed: {}".format(len(all_results)))

    print("\n" + "=" * 70)
    print("OVERVIEW TABLE")
    print("=" * 70)
    print("{:<55} {:>6}".format("File", "Slides"))
    print("-" * 65)
    for r in all_results:
        f = r["file"]
        if len(f) > 54:
            f = "..." + f[-51:]
        print("{:<55} {:>6}".format(f, r["slide_count"]))


if __name__ == "__main__":
    main()
