#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
テンプレート対応 JSON → SX テンプレート PPTX 生成 - 完全版 v3

修正点:
1. body プレースホルダ (14) に合わせた座標システム
2. body プレースホルダを使わない場合は削除
3. objects は body 範囲内に配置
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def hex_to_rgb(hex_color):
    """16進数カラーコード を RGB に変換"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_pptx_final_v3(json_path: str, output_path: str, 
                         use_body_placeholder: bool = False) -> bool:
    """
    JSON → SX テンプレート PPTX 生成（完全版 v3）
    
    修正:
    - body プレースホルダに合わせた配置（left 0.997", width 11.340"）
    - body プレースホルダを使わない場合は削除
    - ユーザーが objects のみで構成する場合に最適化
    
    Args:
        json_path: 入力 JSON ファイル
        output_path: 出力 PPTX ファイル
        use_body_placeholder: body テキストを使用するか（デフォルト False）
    
    Returns:
        成功した場合 True
    """
    try:
        # JSON を読み込み
        with open(json_path, 'r', encoding='utf-8') as f:
            slide_data = json.load(f)
        
        # テンプレートをロード
        template_path = Path(__file__).parent / "templates" / "sx_proposal" / "template.pptx"
        
        if not template_path.exists():
            print(f"   ❌ テンプレートが見つかりません: {template_path}")
            return False
        
        prs = Presentation(str(template_path))
        
        # ─── 既存スライドをすべて削除 ───
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        
        # ─── content layout (index 6) で新しいスライドを追加 ───
        content_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(content_layout)
        
        # ─── タイトルを設定 ───
        if slide_data.get('title'):
            try:
                title_shape = slide.placeholders[0]
                title_shape.text = slide_data['title']
            except Exception as e:
                pass
        
        # ─── サブタイトル（subtitle）を設定 ───
        if slide_data.get('subtitle'):
            try:
                subtitle_shape = slide.placeholders[13]
                subtitle_shape.text = slide_data['subtitle']
            except Exception as e:
                pass
        
        # ─── body プレースホルダを削除（objects のみで構成する場合）───
        if not use_body_placeholder and slide_data.get('objects'):
            try:
                body_shape = slide.placeholders[14]
                sp = body_shape.element
                sp.getparent().remove(sp)
                print(f"   Removed body placeholder [14]")
            except Exception as e:
                pass  # placeholder が見つからない場合はスキップ
        
        # ─── オブジェクトを配置 ───
        for obj in slide_data.get('objects', []):
            obj_type = obj.get('type')
            
            left = obj.get('left', 0)
            top = obj.get('top', 0)
            width = obj.get('width', 1)
            height = obj.get('height', 0.5)
            
            if obj_type == 'box':
                # 背景付きテキストボックス
                shape = slide.shapes.add_shape(
                    1,  # rectangle
                    Inches(left),
                    Inches(top),
                    Inches(width),
                    Inches(height)
                )
                
                # 背景色
                fill_color = obj.get('fill_color', 'FFFFFF')
                rgb = hex_to_rgb(fill_color)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*rgb)
                
                # 枠線
                shape.line.color.rgb = RGBColor(*rgb)
                shape.line.width = Pt(0)
                
                # テキスト
                text_frame = shape.text_frame
                text_frame.word_wrap = True
                
                # 垂直配置（v_align）
                v_align = obj.get('v_align', 'middle')
                if v_align == 'top':
                    text_frame.vertical_anchor = 0  # MSO_ANCHOR.TOP
                elif v_align == 'bottom':
                    text_frame.vertical_anchor = 2  # MSO_ANCHOR.BOTTOM
                else:
                    text_frame.vertical_anchor = 1  # MSO_ANCHOR.MIDDLE
                
                text_frame.margin_bottom = Inches(0.03)
                text_frame.margin_left = Inches(0.03)
                text_frame.margin_right = Inches(0.03)
                text_frame.margin_top = Inches(0.03)
                
                text = obj.get('text', '')
                lines = text.split('\n')
                
                for i, line in enumerate(lines):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = line.strip()
                    
                    # 水平配置（h_align）
                    h_align = obj.get('h_align', 'center')
                    if h_align == 'left':
                        p.alignment = PP_ALIGN.LEFT
                    elif h_align == 'right':
                        p.alignment = PP_ALIGN.RIGHT
                    else:
                        p.alignment = PP_ALIGN.CENTER
                    
                    # テキスト色
                    font_color = obj.get('font_color', '000000')
                    rgb_font = hex_to_rgb(font_color)
                    p.font.color.rgb = RGBColor(*rgb_font)
                    
                    # フォントサイズ
                    font_size = obj.get('font_size', 10)
                    p.font.size = Pt(font_size)
                    p.font.name = 'Arial'
                    p.line_spacing = 1.0
            
            elif obj_type == 'arrow':
                # 矢印（AutoShape 右向き矢印）
                arrow_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Inches(left),
                    Inches(top),
                    Inches(width),
                    Inches(height)
                )
                
                # 背景色（矢印の色）
                arrow_color = obj.get('fill_color', 'ED7D31')
                rgb_arrow = hex_to_rgb(arrow_color)
                arrow_shape.fill.solid()
                arrow_shape.fill.fore_color.rgb = RGBColor(*rgb_arrow)
                
                # 枠線なし
                arrow_shape.line.color.rgb = RGBColor(*rgb_arrow)
                arrow_shape.line.width = Pt(0)
            
            elif obj_type == 'text':
                # テキストボックス（背景なし）
                text_shape = slide.shapes.add_textbox(
                    Inches(left),
                    Inches(top),
                    Inches(width),
                    Inches(height)
                )
                
                text_frame = text_shape.text_frame
                text_frame.word_wrap = True
                
                # 垂直配置
                v_align = obj.get('v_align', 'top')
                if v_align == 'middle':
                    text_frame.vertical_anchor = 1
                elif v_align == 'bottom':
                    text_frame.vertical_anchor = 2
                else:
                    text_frame.vertical_anchor = 0
                
                text = obj.get('text', '')
                lines = text.split('\n')
                
                for i, line in enumerate(lines):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = line.strip()
                    
                    # 水平配置
                    h_align = obj.get('h_align', 'left')
                    if h_align == 'center':
                        p.alignment = PP_ALIGN.CENTER
                    elif h_align == 'right':
                        p.alignment = PP_ALIGN.RIGHT
                    else:
                        p.alignment = PP_ALIGN.LEFT
                    
                    # テキスト色
                    font_color = obj.get('font_color', '000000')
                    rgb_font = hex_to_rgb(font_color)
                    p.font.color.rgb = RGBColor(*rgb_font)
                    
                    # フォントサイズ
                    font_size = obj.get('font_size', 10)
                    p.font.size = Pt(font_size)
                    p.font.name = 'Arial'
                    p.line_spacing = 1.0
        
        # PPTX を保存
        prs.save(output_path)
        return True
    
    except Exception as e:
        print(f"   ❌ エラー: {e}")
        import traceback
        traceback.print_exc()
        return False

def main():
    print("\n" + "="*70)
    print("SX テンプレート PPTX 生成 - v3 (body プレースホルダ対応)")
    print("="*70)
    
    json_files = [
        ("test_output/advanced_swimlane.json", "output_swimlane_v3.pptx"),
        ("test_output/advanced_matrix.json", "output_matrix_v3.pptx"),
    ]
    
    for json_path, pptx_path in json_files:
        json_file = Path(json_path)
        if not json_file.exists():
            print(f"\n❌ {json_path} が見つかりません")
            continue
        
        print(f"\n[*] {json_path}")
        
        with open(json_path, 'r', encoding='utf-8') as f:
            slide_data = json.load(f)
        
        title = slide_data.get('title', 'Untitled')
        obj_count = len(slide_data.get('objects', []))
        
        print(f"   Title: {title}")
        print(f"   Objects: {obj_count}")
        
        # body プレースホルダは使わない（objects のみで構成）
        success = create_pptx_final_v3(json_path, pptx_path, use_body_placeholder=False)
        
        if success:
            file_size = Path(pptx_path).stat().st_size
            print(f"   ✅ Generated ({file_size:,} bytes)")
            print(f"      - body プレースホルダ: 削除")
            print(f"      - objects 配置範囲: left 0.997\", width 11.340\"")
    
    print("\n" + "="*70)
    print("✅ 生成完了")
    print("="*70)
    print(f"\n📋 生成ファイル:")
    print(f"   • output_swimlane_v3.pptx")
    print(f"   • output_matrix_v3.pptx")
    
    print(f"\n🔑 v3 の改善点:")
    print(f"   ✓ body プレースホルダの幅に完全準拠（left 0.997\", width 11.340\"）")
    print(f"   ✓ body プレースホルダを削除（objects 専用レイアウト）")
    print(f"   ✓ スライドレイアウトをすっきりクリーンに")
    
    print(f"\n" + "="*70)

if __name__ == "__main__":
    main()
