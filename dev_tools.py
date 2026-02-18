"""
dev_tools.py - 開発・デバッグ用ユーティリティ

使い方:
  python dev_tools.py env                              # API キー確認
  python dev_tools.py template <path.pptx>             # テンプレートのシェイプ・構造確認
  python dev_tools.py layouts <path.pptx>              # レイアウト別プレースホルダー一覧
  python dev_tools.py layouts                          # デフォルトテンプレートのレイアウト
  python dev_tools.py validate-design <tier2.json>     # スライドJSONのデザインルール検証
"""

import sys
import json
import re
from pathlib import Path


def cmd_env():
    """環境変数（APIキー）の設定確認"""
    from dotenv import load_dotenv
    import os
    load_dotenv()
    ak = os.getenv("ANTHROPIC_API_KEY", "")
    gk = os.getenv("GEMINI_API_KEY", "")
    print("ANTHROPIC_API_KEY:", ak[:20] + "..." if len(ak) > 20 else ak or "(not set)")
    print("GEMINI_API_KEY   :", gk[:20] + "..." if len(gk) > 20 else gk or "(not set)")


def cmd_template(pptx_path: str):
    """テンプレートのスライド・シェイプ構造を表示"""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    print(f"File: {pptx_path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Size: {prs.slide_width.inches:.3f} x {prs.slide_height.inches:.3f} inches")
    print()
    print("=== Layouts ===")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  [{i}] {layout.name}")
    print()
    print("=== Slides ===")
    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        print(f"\nSlide {i+1} [{layout_name}]: {len(slide.shapes)} shapes")
        for s in slide.shapes:
            text = ""
            if s.has_text_frame:
                text = s.text_frame.text[:40].replace("\n", " ")
            try:
                ph_idx = s.placeholder_format.idx if s.is_placeholder else None
            except Exception:
                ph_idx = None
            print(f'  type={s.shape_type} name="{s.name}" ph_idx={ph_idx} text="{text}"')


def cmd_layouts(pptx_path: str):
    """レイアウト別のプレースホルダー一覧を表示"""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    print(f"File: {pptx_path}")
    print()
    print("=== Layout Placeholders ===")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"\n[{i}] {layout.name}")
        for ph in layout.placeholders:
            print(f"  ph_idx={ph.placeholder_format.idx} name=\"{ph.name}\" type={ph.placeholder_format.type}")

    print("\n\n=== Slide Placeholders (actual) ===")
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i+1} [{slide.slide_layout.name}]")
        for ph in slide.placeholders:
            text = ph.text_frame.text[:40].replace("\n", " ") if ph.has_text_frame else ""
            print(f'  ph_idx={ph.placeholder_format.idx} name="{ph.name}" text="{text}"')


def cmd_validate_design(json_path: str):
    """Tier 2 JSON のデザインルール検証"""
    json_path = Path(json_path)
    if not json_path.exists():
        print(f"エラー: ファイルが見つかりません: {json_path}")
        sys.exit(1)
    
    data = json.loads(json_path.read_text(encoding="utf-8"))
    
    # テンプレート特定（プロジェクトの outline.json から読む）
    project_dir = json_path.parent.parent
    outline_path = project_dir / "outline.json"
    template_id = "sx_proposal"  # デフォルト
    if outline_path.exists():
        outline = json.loads(outline_path.read_text(encoding="utf-8"))
        if isinstance(outline, dict):
            template_id = outline.get("template", "sx_proposal")
    
    # design_guide.md を読み込み
    design_guide_path = Path(__file__).parent / "templates" / template_id / "design_guide.md"
    if not design_guide_path.exists():
        print(f"警告: design_guide.md が見つかりません: {design_guide_path}")
        defined_colors = set()
        content_area = {"left_min": 0.5, "left_max": 12.8, "top_min": 1.5, "top_max": 7.0}
    else:
        guide_text = design_guide_path.read_text(encoding="utf-8")
        # 色定義を抽出（カラーコード）
        defined_colors = set(re.findall(r'`([0-9A-Fa-f]{6})`', guide_text))
        # content area を抽出
        content_area = {"left_min": 0.5, "left_max": 12.8, "top_min": 1.5, "top_max": 7.0}
    
    print(f"\n{'='*80}")
    print(f"デザイン検証: {json_path.name}")
    print(f"テンプレート: {template_id}")
    print(f"{'='*80}\n")
    
    errors = []
    warnings = []
    
    objects = data.get("objects", [])
    if not objects:
        print("警告: objects が空です\n")
        return
    
    print(f"オブジェクト数: {len(objects)}\n")
    
    for i, obj in enumerate(objects):
        obj_type = obj.get("type", "不明")
        left = obj.get("left", 0)
        top = obj.get("top", 0)
        width = obj.get("width", 0)
        height = obj.get("height", 0)
        
        # 座標チェック
        if left < content_area["left_min"]:
            errors.append(f"[{i}] {obj_type}: left={left:.2f} < {content_area['left_min']} (左マージン超過)")
        if left + width > content_area["left_max"]:
            errors.append(f"[{i}] {obj_type}: left+width={left+width:.2f} > {content_area['left_max']} (右端超過)")
        if top < content_area["top_min"]:
            warnings.append(f"[{i}] {obj_type}: top={top:.2f} < {content_area['top_min']} (ヘッダー領域に重複の可能性)")
        if top + height > content_area["top_max"]:
            warnings.append(f"[{i}] {obj_type}: top+height={top+height:.2f} > {content_area['top_max']} (フッター領域に重複の可能性)")
        
        # 色チェック
        for color_key in ["fill_color", "font_color"]:
            if color_key in obj:
                color = obj[color_key].upper()
                if color and color not in defined_colors and color not in ("FFFFFF", "000000"):
                    errors.append(f"[{i}] {obj_type}: {color_key}=#{color} は design_guide.md に未定義")
        
        # フォントサイズチェック
        font_size = obj.get("font_size")
        if font_size and (font_size < 9 or font_size > 20):
            warnings.append(f"[{i}] {obj_type}: font_size={font_size} は推奨範囲外 (9-20pt)")
    
    # 結果表示
    if not errors and not warnings:
        print("✅ エラー・警告なし\n")
    else:
        if errors:
            print(f"❌ エラー ({len(errors)}件):\n")
            for err in errors:
                print(f"  {err}")
            print()
        if warnings:
            print(f"⚠️  警告 ({len(warnings)}件):\n")
            for warn in warnings:
                print(f"  {warn}")
            print()
    
    # 定義色一覧
    if defined_colors:
        print(f"定義済み色 ({len(defined_colors)}個):")
        for color in sorted(defined_colors):
            print(f"  #{color}")
        print()


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    cmd = sys.argv[1]
    default_template = str(Path(__file__).parent / "templates" / "sx_proposal" / "template.pptx")

    if cmd == "env":
        cmd_env()
    elif cmd == "template":
        path = sys.argv[2] if len(sys.argv) > 2 else default_template
        cmd_template(path)
    elif cmd == "layouts":
        path = sys.argv[2] if len(sys.argv) > 2 else default_template
        cmd_layouts(path)
    elif cmd == "validate-design":
        if len(sys.argv) < 3:
            print("エラー: JSONファイルパスを指定してください")
            print("使い方: python dev_tools.py validate-design <tier2.json>")
            sys.exit(1)
        cmd_validate_design(sys.argv[2])
    else:
        print(f"Unknown command: {cmd}")
        print(__doc__)
        sys.exit(1)


if __name__ == "__main__":
    main()
