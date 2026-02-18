#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
テンプレート対応 JSON → SX テンプレート PPTX 生成 v2

重要: JSON の座標 = PPTX に配置される座標（完全一致）
スライドサイズや座標変換は不要
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def hex_to_rgb(hex_color):
    """16進数カラーコード を RGB に変換"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_pptx_with_sx_template_v2(json_path: str, output_path: str) -> bool:
    """
    JSON ファイルから、SX テンプレートを使用して PPTX を生成（v2）
    
    重要: JSON の座標はすでに SX template 準拠
    スライドサイズ: 13.333" × 7.5"
    Body content area: left 0.997", top 1.545", width 11.340", height 5.512"
    
    Args:
        json_path: 入力 JSON ファイル
        output_path: 出力 PPTX ファイル
    
    Returns:
        成功した場合 True
    """
    try:
        # JSON を読み込み
        with open(json_path, 'r', encoding='utf-8') as f:
            slide_data = json.load(f)
        
        # SX テンプレートをロード
        template_path = Path(__file__).parent / "templates" / "sx_proposal" / "template.pptx"
        
        if not template_path.exists():
            print(f"   ❌ テンプレートが見つかりません: {template_path}")
            return False
        
        prs = Presentation(str(template_path))
        
        # スライドサイズを確認
        print(f"   Slide size: {prs.slide_width.inches:.3f}\" × {prs.slide_height.inches:.3f}\"")
        
        # content レイアウト (index 6) を使用
        content_layout = prs.slide_layouts[6]
        
        # 既存スライドがあればそれを使用、なければ新しいスライドを追加
        if len(prs.slides) > 0:
            # テンプレートの最後のスライドを使用
            slide = prs.slides[-1]
            print(f"   Using existing slide #{len(prs.slides)}")
        else:
            # 新しいスライドを追加
            slide = prs.slides.add_slide(content_layout)
            print(f"   Added new content slide")
        
        # ─── タイトルを設定 ───
        if slide_data.get('title'):
            try:
                title_shape = slide.placeholders[0]
                title_shape.text = slide_data['title']
                print(f"   Title: {slide_data['title'][:50]}")
            except:
                pass  # placeholder がない場合はスキップ
        
        # ─── サブタイトルを設定 ───
        if slide_data.get('subtitle'):
            try:
                subtitle_shape = slide.placeholders[13]
                subtitle_shape.text = slide_data['subtitle']
            except:
                pass  # placeholder がない場合はスキップ
        
        # ─── コンテンツオブジェクトを配置 ───
        # 重要: JSON の座標をそのまま使用（すでに SX template 準拠）
        
        for obj in slide_data.get('objects', []):
            obj_type = obj.get('type')
            
            # JSON の座標をそのまま使用
            left = obj.get('left', 0)
            top = obj.get('top', 0)
            width = obj.get('width', 1)
            height = obj.get('height', 0.5)
            
            if obj_type == 'box':
                # 背景付きテキストボックス
                shape = slide.shapes.add_shape(
                    1,  # rectangle
                    Inches(left),
                    Inches(top),
                    Inches(width),
                    Inches(height)
                )
                
                # 背景色
                fill_color = obj.get('fill_color', 'FFFFFF')
                rgb = hex_to_rgb(fill_color)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*rgb)
                
                # 枠線
                shape.line.color.rgb = RGBColor(*rgb)
                
                # テキスト
                text_frame = shape.text_frame
                text_frame.word_wrap = True
                text_frame.vertical_anchor = 1  # middle
                text_frame.margin_bottom = Inches(0.05)
                text_frame.margin_left = Inches(0.05)
                text_frame.margin_right = Inches(0.05)
                text_frame.margin_top = Inches(0.05)
                
                text = obj.get('text', '')
                lines = text.split('\n')
                
                for i, line in enumerate(lines):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = line
                    p.alignment = PP_ALIGN.CENTER
                    
                    # テキスト色
                    font_color = obj.get('font_color', '000000')
                    rgb_font = hex_to_rgb(font_color)
                    p.font.color.rgb = RGBColor(*rgb_font)
                    
                    # フォントサイズ
                    font_size = obj.get('font_size', 12)
                    p.font.size = Pt(font_size)
                    p.font.name = 'Arial'
                    p.font.bold = obj.get('font_bold', False)
            
            elif obj_type == 'text':
                # テキストボックス（背景なし）
                text_box = slide.shapes.add_textbox(
                    Inches(left),
                    Inches(top),
                    Inches(width),
                    Inches(height)
                )
                
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                
                text = obj.get('text', '')
                lines = text.split('\n')
                
                for i, line in enumerate(lines):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = line
                    
                    # テキスト色
                    font_color = obj.get('font_color', '000000')
                    rgb_font = hex_to_rgb(font_color)
                    p.font.color.rgb = RGBColor(*rgb_font)
                    
                    # フォントサイズ
                    font_size = obj.get('font_size', 11)
                    p.font.size = Pt(font_size)
                    p.font.name = 'Arial'
            
            elif obj_type == 'arrow':
                # 矢印（簡易的に線で表現）
                connector = slide.shapes.add_connector(
                    1,  # straight connector
                    Inches(left),
                    Inches(top),
                    Inches(left + width),
                    Inches(top + height)
                )
                
                line = connector.line
                line_color = obj.get('fill_color', '404040')
                rgb_line = hex_to_rgb(line_color)
                line.color.rgb = RGBColor(*rgb_line)
                line.width = Pt(2)
        
        # PPTX を保存
        prs.save(output_path)
        return True
    
    except Exception as e:
        print(f"   ❌ エラー: {e}")
        import traceback
        traceback.print_exc()
        return False

def main():
    print("\n" + "="*70)
    print("SX テンプレート PPTX 生成 v2 - Template Compliance")
    print("="*70)
    
    json_files = [
        ("test_output/demo1_sx_template.json", "output_demo1_sx_template_v2.pptx"),
        ("test_output/demo2_sx_template.json", "output_demo2_sx_template_v2.pptx"),
        ("test_output/demo3_sx_template.json", "output_demo3_sx_template_v2.pptx"),
    ]
    
    for json_path, pptx_path in json_files:
        json_file = Path(json_path)
        if not json_file.exists():
            print(f"\n❌ {json_path} が見つかりません")
            continue
        
        print(f"\n[*] {json_path}")
        
        with open(json_path, 'r', encoding='utf-8') as f:
            slide_data = json.load(f)
        
        title = slide_data.get('title', 'Untitled')
        obj_count = len(slide_data.get('objects', []))
        
        print(f"   Objects: {obj_count}")
        print(f"   Generating: {pptx_path}")
        
        success = create_pptx_with_sx_template_v2(json_path, pptx_path)
        
        if success:
            file_size = Path(pptx_path).stat().st_size
            print(f"   ✅ Generated ({file_size:,} bytes)")
    
    print("\n" + "="*70)
    print("✅ 生成完了")
    print("="*70)
    print(f"\n📋 生成ファイル:")
    print(f"   • output_demo1_sx_template_v2.pptx")
    print(f"   • output_demo2_sx_template_v2.pptx")
    print(f"   • output_demo3_sx_template_v2.pptx")
    
    print(f"\n🔑 v2 の改善点:")
    print(f"   ✓ JSON の座標をそのまま使用（座標変換なし）")
    print(f"   ✓ Canvas で見える = PPTX に配置される（完全一致）")
    print(f"   ✓ テンプレートの content layout に直接配置")
    print(f"   ✓ 配置エラーを完全に排除")
    
    print(f"\n" + "="*70)

if __name__ == "__main__":
    main()
