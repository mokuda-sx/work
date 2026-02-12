"""
pptx_engine.py
PPTX生成エンジン（中間言語JSONからPPTXを生成するコアモジュール）

中間言語スキーマ:
[
  {
    "type": "title" | "agenda" | "chapter" | "content" | "end",
    "title": "スライドタイトル",
    "subtitle": "キーメッセージ",
    "body": "本文\n複数行\n箇条書き",
    "objects": [
      {"type": "box",   "text": "ラベル", "left": 0.5, "top": 3.0, "width": 2.5, "height": 0.9,
       "fill_color": "4472C4", "font_color": "FFFFFF", "font_size": 14, "bold": true},
      {"type": "arrow", "left": 3.1, "top": 3.2, "width": 0.5, "height": 0.5, "fill_color": "ED7D31"},
      {"type": "text",  "text": "補足", "left": 1.0, "top": 4.0, "width": 4.0, "height": 0.5,
       "font_color": "404040", "font_size": 11}
    ],
    "images": [
      {
        "prompt": "English image prompt",
        "model": "gemini-3-pro-image-preview",
        "position": "auto",   # テンプレートの「画像挿入位置」シェイプ座標を自動使用
        # または明示的座標指定:
        "left": 7.0, "top": 1.5, "width": 5.5
      }
    ]
  }
]

座標系: 幅13.3 × 高さ7.5インチ（16:9）

テンプレートの「画像挿入位置」シェイプ座標（レイアウト別）:
  title   [0]:  left=1.012, top=2.253, width=11.348, height=4.646
  chapter_photo [4]: left=1.012, top=2.411, width=11.348, height=4.646
  agenda  [2]:  left=6.981, top=0.443, width=5.356, height=6.615
  end     [14]: left=0.997, top=0.831, width=11.348, height=4.646
"""

import os
import io
import json
import ast
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

TEMPLATE_PATH = Path(__file__).parent / "SX_提案書_3.0_16x9.pptx"

LAYOUT = {
    "title":         0,   # ドキュメンテーションタイトルあり（写真あり）
    "chapter":       5,   # チャプタータイトル（写真なし）★シンプル版
    "chapter_photo": 4,   # チャプタータイトル（写真あり）
    "agenda":        2,   # 目次・アジェンダページ（写真あり）
    "content":       6,   # コンテンツページヘッドラインあり
    "end":          14,   # エンドスライド（写真あり）
}

# 「画像挿入位置」シェイプの座標 (layout_index → (left, top, width, height) インチ)
TEMPLATE_IMAGE_AREAS = {
    0:  (1.012, 2.253, 11.348, 4.646),   # title
    4:  (1.012, 2.411, 11.348, 4.646),   # chapter_photo
    2:  (6.981, 0.443,  5.356, 6.615),   # agenda
    14: (0.997, 0.831, 11.348, 4.646),   # end
    16: (0.998, 0.837, 11.348, 4.646),   # end_with_pmark
}

# ─── テンプレート操作 ─────────────────────────────────
def load_template() -> Presentation:
    return Presentation(str(TEMPLATE_PATH))

def remove_all_slides(prs: Presentation):
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    slide_id_list = prs.slides._sldIdLst
    for sld_id in list(slide_id_list):
        rId = sld_id.get(f'{{{r_ns}}}id')
        if rId:
            prs.part.rels.pop(rId)
        slide_id_list.remove(sld_id)

# ─── テキスト設定 ─────────────────────────────────────
def fill_text_frame(tf, text: str):
    """
    テキストフレームにテキストを設定する。
    tf.clear() を使わず XML レベルで操作することで、テンプレートの
    フォント・色・サイズ等のスタイル継承（<a:pPr> / lstStyle）を保持する。
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    lines = [l for l in text.split('\n') if l.strip()]
    if not lines:
        return

    txBody = tf._txBody
    paras  = txBody.findall(qn('a:p'))

    # ── 最初の段落を再利用（<a:pPr> を保持してスタイルを守る） ──
    if paras:
        p0 = paras[0]
        # 既存の run / field を削除（テキスト内容のみ除去、段落書式は残す）
        for child in list(p0.findall(qn('a:r'))) + list(p0.findall(qn('a:fld'))):
            p0.remove(child)
        # 余分な段落を削除
        for p in paras[1:]:
            txBody.remove(p)
    else:
        p0 = etree.SubElement(txBody, qn('a:p'))

    # 最初の行を run として追加（rPr なし → <a:pPr> / lstStyle から継承）
    r = etree.SubElement(p0, qn('a:r'))
    etree.SubElement(r, qn('a:t')).text = lines[0]

    # ── 2行目以降: 新しい段落として追加 ──
    for line in lines[1:]:
        p = etree.SubElement(txBody, qn('a:p'))
        r = etree.SubElement(p, qn('a:r'))
        etree.SubElement(r, qn('a:t')).text = line

def set_placeholder_text(slide, ph_idx: int, text: str) -> bool:
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == ph_idx:
            fill_text_frame(shape.text_frame, text)
            return True
    return False

def set_body_text(slide, text: str):
    for idx in [10, 14, 1, 2]:
        if set_placeholder_text(slide, idx, text):
            return

# ─── オブジェクト設定 ─────────────────────────────────
def parse_objects(objects) -> list[dict]:
    if not objects:
        return []
    if isinstance(objects, list):
        return [o for o in objects if isinstance(o, dict)]
    if isinstance(objects, str):
        try:
            result = ast.literal_eval(objects.strip())
            return result if isinstance(result, list) else []
        except Exception:
            try:
                return json.loads(objects.strip())
            except Exception:
                return []
    return []

def add_objects_to_slide(slide, objects: list[dict]):
    for obj in parse_objects(objects):
        obj_type = obj.get("type", "box")
        left   = Inches(obj.get("left",   1.0))
        top    = Inches(obj.get("top",    2.0))
        width  = Inches(obj.get("width",  2.0))
        height = Inches(obj.get("height", 0.8))

        if obj_type in ("box", "rect"):
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(obj.get("fill_color", "4472C4"))
            shape.line.fill.background()
            text = obj.get("text", "")
            if text:
                tf = shape.text_frame
                tf.word_wrap = True
                tf.clear()
                p = tf.paragraphs[0]
                p.text = text
                p.alignment = PP_ALIGN.CENTER
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size  = Pt(obj.get("font_size", 12))
                run.font.color.rgb = RGBColor.from_string(obj.get("font_color", "FFFFFF"))
                run.font.bold  = obj.get("bold", True)

        elif obj_type == "arrow":
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(obj.get("fill_color", "ED7D31"))
            shape.line.fill.background()

        elif obj_type == "text":
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.text = obj.get("text", "")
            p.alignment = PP_ALIGN.LEFT
            if p.runs:
                run = p.runs[0]
                run.font.size  = Pt(obj.get("font_size", 11))
                run.font.color.rgb = RGBColor.from_string(obj.get("font_color", "404040"))

# ─── 画像生成 (Gemini) ───────────────────────────────
def generate_image_gemini(prompt: str, model: str = "gemini-3-pro-image-preview") -> bytes | None:
    from google import genai
    from google.genai import types as genai_types
    key = os.getenv("GEMINI_API_KEY", "")
    if not key:
        print(f"  [skip] GEMINI_API_KEY 未設定")
        return None
    print(f"  [image] 生成中: {prompt[:50]}...")
    client = genai.Client(api_key=key)
    response = client.models.generate_content(
        model=model,
        contents=prompt,
        config=genai_types.GenerateContentConfig(response_modalities=["IMAGE", "TEXT"]),
    )
    for part in response.candidates[0].content.parts:
        if part.inline_data is not None:
            print(f"  [image] 生成完了")
            return part.inline_data.data
    return None

def _center_crop_to_ratio(img_bytes: bytes, target_w: float, target_h: float) -> bytes:
    """画像を target_w:target_h のアスペクト比で中央クロップする。"""
    from PIL import Image
    img = Image.open(io.BytesIO(img_bytes))
    src_w, src_h = img.size
    target_ratio = target_w / target_h
    src_ratio = src_w / src_h

    if abs(src_ratio - target_ratio) < 0.01:
        return img_bytes  # ほぼ同じ比率ならクロップ不要

    if src_ratio > target_ratio:
        # 元画像が横に広い → 左右をクロップ
        new_w = int(src_h * target_ratio)
        left = (src_w - new_w) // 2
        img = img.crop((left, 0, left + new_w, src_h))
    else:
        # 元画像が縦に長い → 上下をクロップ
        new_h = int(src_w / target_ratio)
        top = (src_h - new_h) // 2
        img = img.crop((0, top, src_w, top + new_h))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def add_images_to_slide(slide, images: list[dict], layout_index: int = -1,
                        slides_dir: Path | None = None):
    """
    images: JSON の images フィールド
    layout_index: テンプレートの「画像挿入位置」座標を参照するためのレイアウトインデックス
    slides_dir: 画像キャッシュの保存先（指定時、生成画像を保存＆次回再利用）
    """
    for img_spec in images:
        # "file" キーがあればローカルファイルを優先使用
        file_path = img_spec.get("file", "")
        if file_path:
            fp = Path(file_path)
            if not fp.is_absolute() and slides_dir:
                fp = slides_dir / fp
            if fp.exists():
                print(f"  [image] キャッシュ使用: {fp.name}")
                img_bytes = fp.read_bytes()
            else:
                print(f"  [image] ファイル未発見: {fp} -> 生成にフォールバック")
                img_bytes = None
                file_path = ""
        else:
            img_bytes = None

        prompt = img_spec.get("prompt", "")
        if not img_bytes and not prompt:
            continue
        model = img_spec.get("model", "gemini-3-pro-image-preview")

        # キャッシュにない場合は Gemini 生成
        if not img_bytes:
            img_bytes = generate_image_gemini(prompt, model=model)
            # 生成成功時、slides_dir があれば自動保存
            if img_bytes and slides_dir:
                cache_dir = slides_dir / "images"
                cache_dir.mkdir(exist_ok=True)
                # ファイル名: prompt先頭20文字をサニタイズ
                safe_name = "".join(c if c.isalnum() or c in "-_ " else "" for c in prompt[:30]).strip()
                cache_path = cache_dir / f"{safe_name}.png"
                cache_path.write_bytes(img_bytes)
                print(f"  [image] 保存: {cache_path.name}")

        if not img_bytes:
            continue

        # position="auto" または left が未指定 → テンプレートの画像エリアを使用
        use_template_area = (
            img_spec.get("position") == "auto"
            or ("left" not in img_spec and layout_index in TEMPLATE_IMAGE_AREAS)
        )

        if use_template_area and layout_index in TEMPLATE_IMAGE_AREAS:
            left_inch, top_inch, width_inch, height_inch = TEMPLATE_IMAGE_AREAS[layout_index]
        else:
            left_inch   = img_spec.get("left",   7.0)
            top_inch    = img_spec.get("top",    1.5)
            width_inch  = img_spec.get("width",  5.5)
            height_inch = img_spec.get("height", None)

        # 配置先の幅・高さが両方分かる場合は中央クロップしてアスペクト比を合わせる
        if height_inch:
            img_bytes = _center_crop_to_ratio(img_bytes, width_inch, height_inch)
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Inches(left_inch), Inches(top_inch),
                Inches(width_inch), Inches(height_inch)
            )
        else:
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Inches(left_inch), Inches(top_inch), Inches(width_inch)
            )

# ─── PNG サムネイル生成 ───────────────────────────────
def export_thumbnails(pptx_path: Path) -> list[Path]:
    """
    PowerShell 経由で PowerPoint COM を呼び出し、各スライドをPNGに書き出す。
    管理者権限不要・追加パッケージ不要（PowerPoint がインストール済みであること）。
    戻り値: 生成した PNG パスのリスト（失敗時は空リスト）
    """
    import subprocess

    thumb_dir = pptx_path.parent / "thumbnails"
    thumb_dir.mkdir(exist_ok=True)
    print(f"  [thumbnail] PNG生成中: {pptx_path.name}")

    pptx_abs  = str(pptx_path.resolve())
    thumb_abs = str(thumb_dir.resolve())

    ps_script = f"""
$ErrorActionPreference = 'Stop'
$ppt = New-Object -ComObject PowerPoint.Application
$ppt.Visible = [Microsoft.Office.Core.MsoTriState]::msoTrue
try {{
    $pres = $ppt.Presentations.Open('{pptx_abs}', $true, $false, $false)
    $pres.Export('{thumb_abs}', 'PNG', 1920, 1080)
    $pres.Close()
    Write-Output "OK:$($pres.Slides.Count)"
}} finally {{
    $ppt.Quit()
    [System.Runtime.Interopservices.Marshal]::ReleaseComObject($ppt) | Out-Null
}}
"""
    try:
        result = subprocess.run(
            ["powershell", "-NoProfile", "-Command", ps_script],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode != 0:
            print(f"  [thumbnail] PowerShell エラー: {result.stderr.strip()[:200]}")
            return []

        png_paths = sorted(thumb_dir.glob("*.PNG")) + sorted(thumb_dir.glob("*.png"))
        for p in png_paths:
            print(f"  [thumbnail] {p.name}")
        print(f"  [thumbnail] 完了: {thumb_dir} ({len(png_paths)}枚)")
        return png_paths

    except subprocess.TimeoutExpired:
        print("  [thumbnail] タイムアウト（120秒）")
        return []
    except Exception as e:
        print(f"  [thumbnail] エラー: {e}")
        return []

# ─── スライド追加 ─────────────────────────────────────
def add_slide(prs: Presentation, slide_data: dict, slides_dir: Path | None = None):
    layout_key   = slide_data.get("type", "content")
    title        = slide_data.get("title",    "")
    subtitle     = slide_data.get("subtitle", "")
    body         = slide_data.get("body",     "")
    objects      = parse_objects(slide_data.get("objects", []))
    images       = slide_data.get("images",   [])

    layout_index = LAYOUT.get(layout_key, LAYOUT["content"])
    layout = prs.slide_layouts[layout_index]
    slide  = prs.slides.add_slide(layout)

    if title:    set_placeholder_text(slide, 0, title)
    if subtitle:
        # title/agenda レイアウトのサブタイトルは ph_idx=1
        # content レイアウトのキーメッセージは ph_idx=13
        subtitle_idx = 1 if layout_key in ("title", "agenda") else 13
        set_placeholder_text(slide, subtitle_idx, subtitle)
    if body:     set_body_text(slide, body)
    if objects:  add_objects_to_slide(slide, objects)
    if images:   add_images_to_slide(slide, images, layout_index=layout_index,
                                     slides_dir=slides_dir)

    return slide

# ─── スライドディレクトリから結合 ──────────────────────
def build_from_slides_dir(slides_dir: Path, output_path: Path,
                          export_png: bool = False) -> Path:
    """
    slides_dir 内の NN_*.json を番号順に読み込んで PPTX を組み立てる（Tier 2 結合）。
    ファイル名の先頭数字でソートするため、00_title.json → 01_agenda.json の順が保証される。
    Returns: 保存したファイルのPath
    """
    slides_dir = Path(slides_dir)
    json_files = sorted(slides_dir.glob("*.json"), key=lambda p: p.name)
    if not json_files:
        raise FileNotFoundError(f"スライドファイルが見つかりません: {slides_dir}")
    outline = []
    for f in json_files:
        slide_data = json.loads(f.read_text(encoding="utf-8"))
        if isinstance(slide_data, dict):
            outline.append(slide_data)
    print(f"  {len(outline)}枚のスライドを読み込み: {slides_dir}")
    return build_pptx(outline, output_path, export_png=export_png,
                      slides_dir=slides_dir)

# ─── メイン生成関数 ───────────────────────────────────
def build_pptx(outline: list[dict], output_path: str | Path,
               export_png: bool = False,
               slides_dir: Path | None = None) -> Path:
    """
    outline: 中間言語JSONリスト
    output_path: 出力先パス
    export_png: True の場合 PowerPoint COM で PNG サムネイルも生成
    slides_dir: 画像キャッシュの保存/参照先
    Returns: 保存したファイルのPath
    """
    prs = load_template()
    remove_all_slides(prs)
    for i, slide_data in enumerate(outline):
        if not isinstance(slide_data, dict):
            continue
        slide_type = slide_data.get("type", "content")
        print(f"  [{i+1}/{len(outline)}] {slide_type}: {slide_data.get('title', '')[:30]}")
        add_slide(prs, slide_data, slides_dir=slides_dir)
    output_path = Path(output_path)
    prs.save(str(output_path))

    if export_png:
        export_thumbnails(output_path)

    return output_path
