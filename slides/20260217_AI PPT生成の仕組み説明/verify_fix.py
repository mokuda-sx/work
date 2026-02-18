from pptx import Presentation
from pptx.util import Pt

prs = Presentation('20260217_1802_fontsize_fix.pptx')
print(f"Total slides: {len(prs.slides)}")

# Try all slides
for slide_idx, slide in enumerate(prs.slides):
    print(f"\n=== Slide {slide_idx} ({slide.shapes.__len__()} shapes) ===\n")
    
    for i, shape in enumerate(slide.shapes):
        if hasattr(shape, 'text_frame'):
            text = shape.text
            if text.strip() and any(x in text for x in ['構成方針', 'を決定', 'AIが', 'パターン', 'アウトライン', 'Tier']):
                print(f"Box {i}: {repr(text[:50])}")
                frame = shape.text_frame
                
                for p_idx, paragraph in enumerate(frame.paragraphs):
                    for r_idx, run in enumerate(paragraph.runs):
                        font_size = run.font.size
                        if font_size:
                            print(f"  [{p_idx}:{r_idx}] {repr(run.text)} | {font_size.pt}pt ✓")
                        else:
                            print(f"  [{p_idx}:{r_idx}] {repr(run.text)} | None ✗")

