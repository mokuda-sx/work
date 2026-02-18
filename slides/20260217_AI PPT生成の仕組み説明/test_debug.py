import sys
import json
from pathlib import Path

work_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(work_root))

from pptx_engine import build_from_slides_dir, add_slide, get_template_config, load_template, remove_all_slides
from pptx import Presentation

# Load the JSON directly
slides_dir = Path(__file__).parent / "slides"
slide_data = json.loads((slides_dir / "01_content.json").read_text(encoding='utf-8'))

print(f"Loaded slide data with {len(slide_data.get('objects', []))} objects")

# Check object types
objects = slide_data.get('objects', [])
for i, obj in enumerate(objects):
    print(f"  {i}: type={obj.get('type')}")

# Now try to build the PPTX manually to debug
config = get_template_config("sx_proposal")
prs = load_template(config)
remove_all_slides(prs)

print(f"\nAdding slide...")
add_slide(prs, slide_data, config=config, slides_dir=slides_dir)

output_path = Path(__file__).parent / "test_debug.pptx"
prs.save(str(output_path))

# Now check what was created
prs_check = Presentation(str(output_path))
slide = prs_check.slides[0]
print(f"\nCreated slide with {len(slide.shapes)} shapes:")
for i, shape in enumerate(slide.shapes):
    shape_name = shape.name if hasattr(shape, 'name') else 'Unknown'
    is_ascii = shape_name.isascii() if isinstance(shape_name, str) else False
    
    if hasattr(shape, 'text_frame') and shape.text.strip():
        print(f"  {i}: TEXT | {shape.text[:30]}")
    elif "Arrow" in shape_name:
        print(f"  {i}: ARROW | {shape_name}")
    elif hasattr(shape, 'text_frame'):
        print(f"  {i}: EMPTY_TEXTBOX | {shape_name}")
    else:
        print(f"  {i}: SHAPE | {shape_name}")

print(f"\nDone: {output_path}")
