from pptx import Presentation

prs = Presentation('20260217_1802_fontsize_fix.pptx')
slide = prs.slides[1]

print('=== Slide 1 全オブジェクト監査 ===\n')

text_count = 0
arrow_count = 0
rect_count = 0

for i, shape in enumerate(slide.shapes):
    shape_name = shape.name if hasattr(shape, 'name') else 'Unknown'
    
    # Arrow detection
    if "Arrow" in shape_name:
        arrow_count += 1
        print(f'{i}: ARROW | {shape_name}')
        print(f'   pos: ({shape.left.inches:.2f}", {shape.top.inches:.2f}") size: {shape.width.inches:.2f}"×{shape.height.inches:.2f}"')
        try:
            if hasattr(shape, 'fill'):
                rgb = shape.fill.fore_color.rgb
                print(f'   fill: #{rgb}')
        except:
            pass
    # Text detection (non-arrow)
    elif hasattr(shape, 'text_frame'):
        text = shape.text.strip()
        if text:
            text_count += 1
            print(f'{i}: BOX/TEXT | {repr(text[:50])}')
            print(f'   pos: ({shape.left.inches:.2f}", {shape.top.inches:.2f}") size: {shape.width.inches:.2f}"×{shape.height.inches:.2f}"')
            
            # 背景色
            try:
                if hasattr(shape, 'fill'):
                    rgb = shape.fill.fore_color.rgb
                    print(f'   fill: #{rgb}')
            except:
                pass
            
            # フォント情報
            frame = shape.text_frame
            for p_idx, p in enumerate(frame.paragraphs):
                for r_idx, r in enumerate(p.runs):
                    sz = r.font.size
                    print(f'      [{p_idx}:{r_idx}] {repr(r.text[:30])} | {sz.pt if sz else "None"}pt')
    print()

print(f'\n=== Summary ===')
print(f'Text boxes: {text_count}')
print(f'Arrows: {arrow_count}')
print(f'Total shapes: {len(slide.shapes)}')


