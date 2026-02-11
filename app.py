"""
AI PowerPoint Generator
会社テンプレートをベースにAIと対話しながらPPTXを生成するツール
"""

import os
import io
import json
from pathlib import Path
from lxml import etree

import streamlit as st
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import anthropic
from google import genai
from google.genai import types as genai_types

load_dotenv()

# ─── 定数 ────────────────────────────────────────────
TEMPLATE_PATH = Path(r"C:\ProgramData\test\work\SX_提案書_3.0_16x9.pptx")
OUTPUT_DIR = Path(r"C:\ProgramData\test\work\output")
OUTPUT_DIR.mkdir(exist_ok=True)

# テンプレートのレイアウト番号
LAYOUT = {
    "title":    0,   # ドキュメンテーションタイトル
    "chapter":  4,   # チャプタータイトル
    "agenda":   2,   # 詳細・アジェンダページ
    "content":  6,   # コンテンツページ（ヘッドラインあり）
    "end":     14,   # エンドスライド
}

# ─── API クライアント初期化 ───────────────────────────
def get_claude_client():
    key = os.getenv("ANTHROPIC_API_KEY", "")
    if not key or key == "your_anthropic_api_key_here":
        return None
    return anthropic.Anthropic(api_key=key)

# ─── AI 会話 (Claude) ────────────────────────────────
def chat_with_claude(messages: list[dict], system_prompt: str = "") -> str:
    client = get_claude_client()
    if client is None:
        return "⚠️ ANTHROPIC_API_KEY が設定されていません。.env ファイルを確認してください。"
    response = client.messages.create(
        model="claude-sonnet-4-5-20250929",
        max_tokens=8192,
        system=system_prompt,
        messages=messages,
    )
    return response.content[0].text

# ─── 画像生成 (Gemini) ───────────────────────────────
def generate_image_gemini(prompt: str, model: str = "gemini-3-pro-image-preview") -> bytes | None:
    key = os.getenv("GEMINI_API_KEY", "")
    if not key or key == "your_gemini_api_key_here":
        st.warning("⚠️ GEMINI_API_KEY が設定されていません。")
        return None
    client = genai.Client(api_key=key)
    response = client.models.generate_content(
        model=model,
        contents=prompt,
        config=genai_types.GenerateContentConfig(
            response_modalities=["IMAGE", "TEXT"]
        ),
    )
    for part in response.candidates[0].content.parts:
        if part.inline_data is not None:
            return part.inline_data.data
    return None

# ─── PPTX 操作 ───────────────────────────────────────
def load_template() -> Presentation:
    return Presentation(str(TEMPLATE_PATH))

def remove_all_slides(prs: Presentation):
    """テンプレートの既存スライドを全削除（OPCパーツも含む・レイアウト・テーマは保持）"""
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    slide_id_list = prs.slides._sldIdLst
    for sld_id in list(slide_id_list):
        rId = sld_id.get(f'{{{r_ns}}}id')
        if rId:
            prs.part.rels.pop(rId)   # OPCリレーション＆パーツも削除
        slide_id_list.remove(sld_id)

def fill_text_frame(tf, text: str):
    """テキストフレームに複数行テキストを設定（書式は既存スタイル継承）"""
    tf.clear()
    lines = [l for l in text.split('\n') if l.strip()]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line

def set_placeholder_text(slide, ph_idx: int, text: str) -> bool:
    """指定ph_idxのプレースホルダーにテキストをセット。成功したらTrue"""
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == ph_idx:
            fill_text_frame(shape.text_frame, text)
            return True
    return False

def set_body_text(slide, text: str):
    """bodyテキストをセット。ph_idx=10→14→1の順で試みる"""
    for idx in [10, 14, 1, 2]:
        if set_placeholder_text(slide, idx, text):
            return

def parse_objects(objects) -> list[dict]:
    """objects フィールドが文字列や None の場合もリストに変換"""
    if not objects:
        return []
    if isinstance(objects, list):
        return [o for o in objects if isinstance(o, dict)]
    if isinstance(objects, str):
        import ast
        try:
            result = ast.literal_eval(objects.strip())
            return result if isinstance(result, list) else []
        except Exception:
            try:
                result = json.loads(objects.strip())
                return result if isinstance(result, list) else []
            except Exception:
                return []
    return []

def add_objects_to_slide(slide, objects):
    """
    objects: [
      {"type": "box", "text": "テキスト", "left": 1.0, "top": 2.0, "width": 2.0, "height": 0.8,
       "fill_color": "4472C4", "font_color": "FFFFFF", "font_size": 14},
      {"type": "arrow", "left": 3.1, "top": 2.35, "width": 0.4, "height": 0.3},
    ]
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    for obj in objects:
        obj_type = obj.get("type", "box")
        left   = Inches(obj.get("left", 1.0))
        top    = Inches(obj.get("top", 2.0))
        width  = Inches(obj.get("width", 2.0))
        height = Inches(obj.get("height", 0.8))

        if obj_type in ("box", "rect"):
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                left, top, width, height
            )
            # 塗り色
            fill_hex = obj.get("fill_color", "4472C4")
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
            # 枠線なし
            shape.line.fill.background()
            # テキスト
            text = obj.get("text", "")
            if text:
                tf = shape.text_frame
                tf.word_wrap = True
                tf.clear()
                p = tf.paragraphs[0]
                p.text = text
                p.alignment = PP_ALIGN.CENTER
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(obj.get("font_size", 12))
                font_hex = obj.get("font_color", "FFFFFF")
                run.font.color.rgb = RGBColor.from_string(font_hex)
                run.font.bold = obj.get("bold", True)

        elif obj_type == "arrow":
            # 右矢印
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                left, top, width, height
            )
            arrow_hex = obj.get("fill_color", "595959")
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(arrow_hex)
            shape.line.fill.background()

        elif obj_type == "text":
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.text = obj.get("text", "")
            p.alignment = PP_ALIGN.LEFT
            if p.runs:
                run = p.runs[0]
                run.font.size = Pt(obj.get("font_size", 11))
                font_hex = obj.get("font_color", "000000")
                run.font.color.rgb = RGBColor.from_string(font_hex)

def add_slide(prs: Presentation, slide_data: dict) -> object:
    """slide_dataからスライドを生成"""
    layout_key  = slide_data.get("type", "content")
    title       = slide_data.get("title", "")
    body        = slide_data.get("body", "")
    subtitle    = slide_data.get("subtitle", "")
    objects     = parse_objects(slide_data.get("objects", []))

    layout_idx = LAYOUT.get(layout_key, LAYOUT["content"])
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    if title:
        set_placeholder_text(slide, 0, title)
    if subtitle:
        set_placeholder_text(slide, 13, subtitle)
    if body:
        set_body_text(slide, body)
    if objects:
        add_objects_to_slide(slide, objects)

    return slide

def insert_image_to_slide(slide, image_bytes: bytes, left_inch=1.0, top_inch=2.0, width_inch=5.0):
    image_stream = io.BytesIO(image_bytes)
    slide.shapes.add_picture(image_stream, Inches(left_inch), Inches(top_inch), Inches(width_inch))

def build_pptx_from_outline(outline: list[dict]) -> bytes:
    prs = load_template()
    remove_all_slides(prs)          # テンプレートの既存スライドを削除
    for slide_data in outline:
        add_slide(prs, slide_data)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ─── アウトライン生成プロンプト ──────────────────────
OUTLINE_SYSTEM_PROMPT = """あなたはプロのコンサルタントです。
ユーザーの要望を聞いて、提案書のアウトラインをJSON形式で生成してください。

出力形式（JSON配列のみ。説明文・マークダウン記法は一切不要）:
[
  {
    "type": "title",
    "title": "提案書タイトル（20〜35文字）",
    "subtitle": "2026年X月　クライアント名御中"
  },
  {
    "type": "agenda",
    "title": "目次",
    "body": "1. 背景と課題\n2. 提案内容\n3. 期待効果\n4. 実施スケジュール"
  },
  {
    "type": "chapter",
    "title": "1. 背景と課題（20〜30文字）"
  },
  {
    "type": "content",
    "title": "スライドタイトル（20〜35文字）",
    "subtitle": "キーメッセージ：このスライドで伝えたい1文（40〜70文字）",
    "body": "・箇条書き1行目（30〜50文字）\n・箇条書き2行目（30〜50文字）\n・箇条書き3行目（30〜50文字）\n・箇条書き4行目（30〜50文字）",
    "objects": []
  },
  {
    "type": "end"
  }
]

各フィールドの文字数ガイド：
- title（タイトル）: 20〜35文字
- subtitle（キーメッセージ）: 40〜70文字
- body（本文）: 3〜5行の箇条書き、各行30〜50文字、行間は\\nで区切る
- objects: 図解が有効なスライドのみ使用（空配列でもよい）

objectsフィールドで簡単な図解を追加できます（省略可）:
[
  {"type": "box", "text": "現状", "left": 0.5, "top": 3.0, "width": 2.5, "height": 0.9, "fill_color": "4472C4", "font_color": "FFFFFF", "font_size": 14},
  {"type": "arrow", "left": 3.1, "top": 3.2, "width": 0.5, "height": 0.5, "fill_color": "595959"},
  {"type": "box", "text": "提案後", "left": 3.7, "top": 3.0, "width": 2.5, "height": 0.9, "fill_color": "ED7D31", "font_color": "FFFFFF", "font_size": 14}
]

スライドの座標系: 幅13.3インチ × 高さ7.5インチ
- タイトルエリア: top 0〜1.5インチ
- コンテンツエリア: top 1.5〜7.0インチ

typeの種類:
- title: 表紙（titleとsubtitleのみ）
- agenda: 目次（bodyに番号付きリスト）
- chapter: 章区切り（titleのみ）
- content: 通常コンテンツページ（title + subtitle + body + objects）
- end: 最終ページ（空でよい）

必ずJSON配列のみを返してください。```json ``` などのマークダウン記法は使わないこと。
"""

def generate_outline_from_ai(user_request: str, conversation_history: list) -> str:
    messages = conversation_history + [{"role": "user", "content": user_request}]
    return chat_with_claude(messages, system_prompt=OUTLINE_SYSTEM_PROMPT)

# ─── Streamlit UI ────────────────────────────────────
def main():
    st.set_page_config(page_title="AI PPTX Generator", page_icon="📊", layout="wide")
    st.title("📊 AI PowerPoint Generator")
    st.caption("SXテンプレートをベースにAIと対話しながら提案書を作成")

    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "outline" not in st.session_state:
        st.session_state.outline = []
    if "generated_image" not in st.session_state:
        st.session_state.generated_image = None

    # ─── サイドバー ───────────────────────────────────
    with st.sidebar:
        st.header("⚙️ 設定")
        anthropic_key = st.text_input("Anthropic API Key", value=os.getenv("ANTHROPIC_API_KEY", ""), type="password")
        gemini_key    = st.text_input("Gemini API Key",    value=os.getenv("GEMINI_API_KEY", ""),    type="password")

        if st.button("APIキーを保存"):
            env_path = Path(r"C:\ProgramData\test\work\.env")
            env_path.write_text(
                f'ANTHROPIC_API_KEY="{anthropic_key}"\nGEMINI_API_KEY="{gemini_key}"\n',
                encoding="utf-8"
            )
            os.environ["ANTHROPIC_API_KEY"] = anthropic_key
            os.environ["GEMINI_API_KEY"]    = gemini_key
            st.success("保存しました")

        st.divider()
        st.header("🖼️ 画像生成")
        image_model = st.selectbox(
            "モデル",
            ["gemini-3-pro-image-preview", "gemini-2.5-flash-image"],
            help="Pro: 高品質 / Flash: 高速"
        )
        image_prompt = st.text_area(
            "画像プロンプト（英語推奨）", height=100,
            placeholder="e.g. A modern factory with digital automation, clean professional illustration"
        )
        target_slide = st.number_input("挿入先スライド番号（1始まり）", min_value=1, value=1, step=1)
        img_left  = st.slider("左位置 (inch)", 0.0, 12.0, 7.0, 0.1)
        img_top   = st.slider("上位置 (inch)", 0.0, 6.5,  1.5, 0.1)
        img_width = st.slider("幅 (inch)",     1.0, 10.0, 5.5, 0.1)

        if st.button("Geminiで画像生成"):
            with st.spinner("画像生成中..."):
                try:
                    img_bytes = generate_image_gemini(image_prompt, model=image_model)
                    if img_bytes:
                        st.session_state.generated_image = img_bytes
                        st.session_state.generated_image_params = {
                            "slide": int(target_slide) - 1,
                            "left": img_left, "top": img_top, "width": img_width
                        }
                        st.image(img_bytes, caption="生成された画像")
                        st.success("生成完了！PPTXを生成するとこの画像が挿入されます。")
                    else:
                        st.error("画像が返されませんでした")
                except Exception as e:
                    st.error(f"エラー: {e}")

    # ─── メインエリア ─────────────────────────────────
    col1, col2 = st.columns([1, 1])

    with col1:
        st.header("💬 AIとの対話")

        for msg in st.session_state.messages:
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])

        if user_input := st.chat_input("提案書の内容を教えてください"):
            st.session_state.messages.append({"role": "user", "content": user_input})
            with st.chat_message("user"):
                st.markdown(user_input)

            with st.chat_message("assistant"):
                with st.spinner("AIが考えています..."):
                    outline_keywords = ["アウトライン", "作って", "生成", "提案書", "スライド", "まとめて", "構成"]
                    if any(kw in user_input for kw in outline_keywords):
                        response = generate_outline_from_ai(user_input, st.session_state.messages[:-1])
                        # JSON抽出（```json ... ``` ブロックにも対応）
                        raw = response.strip()
                        if raw.startswith("```"):
                            raw = raw.split("```")[1]
                            if raw.startswith("json"):
                                raw = raw[4:]
                            raw = raw.strip()
                        try:
                            outline = json.loads(raw)
                            st.session_state.outline = outline
                            display = (
                                f"✅ アウトラインを生成しました（{len(outline)}スライド）。"
                                f"右側で確認・編集してください。\n\n"
                                f"```json\n{json.dumps(outline, ensure_ascii=False, indent=2)}\n```"
                            )
                        except json.JSONDecodeError:
                            display = f"⚠️ JSONの解析に失敗しました。もう一度試してください。\n\n```\n{response}\n```"
                    else:
                        display = chat_with_claude(
                            st.session_state.messages,
                            system_prompt="あなたはプレゼンテーション作成の専門家です。日本語で簡潔に回答してください。"
                        )

                st.markdown(display)
                st.session_state.messages.append({"role": "assistant", "content": display})

    with col2:
        st.header("📋 アウトライン編集 & PPTX生成")

        if st.session_state.outline:
            st.caption(f"📑 {len(st.session_state.outline)} スライド")
            edited_outline = st.data_editor(
                st.session_state.outline,
                num_rows="dynamic",
                use_container_width=True,
                column_config={
                    "type":     st.column_config.SelectboxColumn("タイプ",         options=["title","agenda","chapter","content","end"]),
                    "title":    st.column_config.TextColumn("タイトル",             width="medium"),
                    "subtitle": st.column_config.TextColumn("キーメッセージ",       width="medium"),
                    "body":     st.column_config.TextColumn("本文（\\nで改行）",    width="large"),
                    "objects":  st.column_config.Column("図解オブジェクト",         width="small"),
                }
            )
            st.session_state.outline = edited_outline

            st.divider()
            col_a, col_b = st.columns(2)
            with col_a:
                output_filename = st.text_input("ファイル名", value="提案書_draft.pptx")
            with col_b:
                st.write("")
                st.write("")
                if st.button("🚀 PPTXを生成・ダウンロード", type="primary"):
                    with st.spinner("PPTX生成中..."):
                        try:
                            pptx_bytes_io = io.BytesIO()
                            prs = load_template()
                            remove_all_slides(prs)
                            for slide_data in st.session_state.outline:
                                if not isinstance(slide_data, dict):
                                    continue
                                s = add_slide(prs, slide_data)
                            # 生成画像があれば指定スライドに挿入
                            img = st.session_state.get("generated_image")
                            img_params = st.session_state.get("generated_image_params", {})
                            if img and img_params:
                                idx = img_params.get("slide", 0)
                                if 0 <= idx < len(prs.slides):
                                    insert_image_to_slide(
                                        prs.slides[idx], img,
                                        left_inch=img_params["left"],
                                        top_inch=img_params["top"],
                                        width_inch=img_params["width"]
                                    )
                            prs.save(pptx_bytes_io)
                            pptx_bytes_io.seek(0)
                            st.download_button(
                                label="⬇️ ダウンロード",
                                data=pptx_bytes_io.read(),
                                file_name=output_filename,
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            )
                            st.success(f"✅ 生成完了！（{len(st.session_state.outline)}スライド）")
                        except Exception as e:
                            st.error(f"エラー: {e}")
                            import traceback
                            st.code(traceback.format_exc())
        else:
            st.info("左のチャットで提案書の内容を入力するとアウトラインが生成されます。")
            if st.button("📄 サンプルアウトラインを読み込む"):
                st.session_state.outline = [
                    {"type": "title",   "title": "DX推進による製造業務効率化提案",      "subtitle": "2026年2月　〇〇株式会社御中", "body": "", "objects": []},
                    {"type": "agenda",  "title": "目次", "subtitle": "",                 "body": "1. 現状と課題\n2. 提案内容\n3. 期待効果\n4. 実施スケジュール", "objects": []},
                    {"type": "chapter", "title": "1. 現状と課題",                        "subtitle": "", "body": "", "objects": []},
                    {"type": "content", "title": "現状の業務課題",
                     "subtitle": "紙業務が主体で月200時間以上の非効率が発生している",
                     "body": "・申請業務の60%が紙ベースで運用されており処理に時間がかかる\n・データ入力の重複作業が月200時間発生しヒューマンエラーも多い\n・情報共有がメール・電話中心でリアルタイム性に欠ける\n・現場と管理部門の情報ギャップが意思決定スピードを低下させている",
                     "objects": [
                         {"type": "box",   "text": "現状\n紙業務60%", "left": 0.5, "top": 5.5, "width": 2.8, "height": 0.9, "fill_color": "C00000", "font_color": "FFFFFF", "font_size": 13},
                         {"type": "arrow", "left": 3.4,               "top": 5.7,  "width": 0.6, "height": 0.5, "fill_color": "ED7D31"},
                         {"type": "box",   "text": "目標\nDX化100%",  "left": 4.1, "top": 5.5, "width": 2.8, "height": 0.9, "fill_color": "4472C4", "font_color": "FFFFFF", "font_size": 13},
                     ]},
                    {"type": "content", "title": "RPA・ペーパーレス化による改善提案",
                     "subtitle": "段階的な自動化により3年でROI 300%・月150時間の削減を実現",
                     "body": "・フェーズ1（3ヶ月）: 申請書類の電子化・ワークフロー導入\n・フェーズ2（6ヶ月）: RPAによるデータ入力・転記作業の自動化\n・フェーズ3（12ヶ月）: BI連携によるリアルタイムダッシュボード整備\n・期待効果: 月150時間削減・エラー率90%低減・コスト年間2,400万円削減",
                     "objects": []},
                    {"type": "end",     "title": "", "subtitle": "", "body": "", "objects": []},
                ]
                st.rerun()

if __name__ == "__main__":
    main()
