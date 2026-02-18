from pptx import Presentation
from pptx.util import Pt

prs = Presentation('slides/20260217_AI PPT生成の仕組み説明/20260217_1735_readability_fix.pptx')
slide = prs.slides[1]

print("=== Slide 1 - テキストボックスのフォントサイズ検証 ===\n")

for i, shape in enumerate(slide.shapes):
    if hasattr(shape, 'text_frame'):
        text = shape.text
        if text.strip():
            frame = shape.text_frame
            print(f"Box {i}: {repr(text[:40])}")
            print(f"  Position: {shape.left}, {shape.top}")
            print(f"  Size: {shape.width} × {shape.height}")
            
            for p_idx, paragraph in enumerate(frame.paragraphs):
                print(f"  Paragraph {p_idx}:")
                for r_idx, run in enumerate(paragraph.runs):
                    font_size = run.font.size
                    if font_size:
                        print(f"    Run {r_idx}: {repr(run.text)} | font_size: {font_size.pt}pt")
                    else:
                        print(f"    Run {r_idx}: {repr(run.text)} | font_size: None (inherited)")
            print()
