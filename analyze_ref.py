#!/usr/bin/env python3
"""
analyze_ref.py -- 参照 PPTX を分析して refs/ ライブラリに登録するツール

使い方:
    python analyze_ref.py refs/sx/filename.pptx --id sx_ctc_jre_suica
    python analyze_ref.py refs/jr/filename.pptx --id jr_20251209_teireikai
    python analyze_ref.py --list         # 登録済み一覧

出力:
    refs/<family>/<id>/analysis.json   -- 構造分析（機械可読）
    refs/<family>/<id>/thumbnails/     -- PNG サムネイル
    refs/index.json                    -- 全参照作品の目次
"""

import argparse
import json
import re
import shutil
import sys
import tempfile
from collections import Counter
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

REFS_DIR = Path(__file__).parent / "refs"
INDEX_PATH = REFS_DIR / "index.json"


# ─── index.json 管理 ─────────────────────────────────────

def load_index() -> dict:
    if INDEX_PATH.exists():
        return json.loads(INDEX_PATH.read_text(encoding="utf-8"))
    return {"refs": []}


def save_index(index: dict):
    INDEX_PATH.write_text(
        json.dumps(index, ensure_ascii=False, indent=2),
        encoding="utf-8"
    )


# ─── 形状分析 ────────────────────────────────────────────

def rgb_to_hex(color) -> str | None:
    try:
        return f"{color.rgb:06X}"
    except Exception:
        return None


def analyze_shape(shape) -> dict:
    info = {
        "name": shape.name,
        "shape_type": str(shape.shape_type).split(".")[-1],
        "left":   round(shape.left   / 914400, 3),
        "top":    round(shape.top    / 914400, 3),
        "width":  round(shape.width  / 914400, 3),
        "height": round(shape.height / 914400, 3),
    }

    # テキスト
    if shape.has_text_frame:
        texts = []
        font_sizes = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                t = run.text.strip()
                if t:
                    texts.append(t)
                if run.font.size:
                    font_sizes.append(round(run.font.size / 12700))
        if texts:
            info["text"] = " / ".join(texts[:3])
        if font_sizes:
            info["font_sizes"] = sorted(set(font_sizes))

    # 塗り色
    try:
        fill = shape.fill
        if fill.type is not None:
            color = rgb_to_hex(fill.fore_color)
            if color:
                info["fill_color"] = color
    except Exception:
        pass

    return info


def infer_slide_type(slide) -> str:
    """レイアウト名からスライド種別を推定"""
    name = slide.slide_layout.name.lower()
    if "title" in name and "content" not in name:
        return "title"
    if any(w in name for w in ("section", "chapter", "divider")):
        return "chapter"
    if "blank" in name:
        return "blank"
    return "content"


# ─── PPTX 分析 ───────────────────────────────────────────

def analyze_pptx(pptx_path: Path) -> dict:
    prs = Presentation(str(pptx_path))

    slides_data = []
    all_colors = []
    all_font_sizes = []

    for i, slide in enumerate(prs.slides):
        shapes_data = []
        for shape in slide.shapes:
            s = analyze_shape(shape)
            shapes_data.append(s)
            if "fill_color" in s:
                all_colors.append(s["fill_color"])
            for fs in s.get("font_sizes", []):
                all_font_sizes.append(fs)

        title_text = ""
        try:
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()[:60]
        except Exception:
            pass

        slides_data.append({
            "index": i,
            "layout_name": slide.slide_layout.name,
            "inferred_type": infer_slide_type(slide),
            "title": title_text,
            "shape_count": len(shapes_data),
            "shapes": shapes_data,
        })

    color_counts = Counter(
        c for c in all_colors
        if c and c.upper() not in ("FFFFFF", "000000", "FEFEFE")
    )
    font_size_counts = Counter(all_font_sizes)

    width_in = round(prs.slide_width  / 914400, 2)
    height_in = round(prs.slide_height / 914400, 2)

    return {
        "slide_count": len(prs.slides),
        "slide_size":  {"width": width_in, "height": height_in},
        "design_observations": {
            "dominant_colors":     [c for c, _ in color_counts.most_common(8)],
            "font_sizes_observed": sorted(font_size_counts.keys()),
        },
        "slides": slides_data,
    }


# ─── サムネイル生成 ──────────────────────────────────────

def generate_thumbnails(pptx_path: Path, thumb_dir: Path) -> list[Path]:
    """
    pptx_path を一時ディレクトリにコピーして export_thumbnails を呼び出し、
    生成した PNG を thumb_dir に移動する。
    """
    sys.path.insert(0, str(Path(__file__).parent))
    from pptx_engine import export_thumbnails

    thumb_dir.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory() as tmp:
        tmp_pptx = Path(tmp) / pptx_path.name
        shutil.copy2(pptx_path, tmp_pptx)

        pngs = export_thumbnails(tmp_pptx)

        copied = []
        for png in pngs:
            dest = thumb_dir / png.name
            shutil.copy2(png, dest)
            copied.append(dest)

    return copied


# ─── 登録 ────────────────────────────────────────────────

def register(pptx_path: Path, ref_id: str, family: str,
             skip_thumbnail: bool = False) -> dict:
    """PPTX を分析して refs/ ライブラリに登録する"""

    out_dir = REFS_DIR / family / ref_id
    out_dir.mkdir(parents=True, exist_ok=True)

    print(f"[analyze] {pptx_path.name}")
    analysis = analyze_pptx(pptx_path)
    analysis.update({
        "id":           ref_id,
        "family":       family,
        "source_file":  str(pptx_path),
        "analyzed_at":  datetime.now().strftime("%Y-%m-%d"),
    })

    # analysis.json
    analysis_path = out_dir / "analysis.json"
    analysis_path.write_text(
        json.dumps(analysis, ensure_ascii=False, indent=2),
        encoding="utf-8"
    )
    print(f"  analysis.json: {analysis_path}")

    # サムネイル
    thumb_dir = out_dir / "thumbnails"
    png_count = 0
    if not skip_thumbnail:
        print("  サムネイル生成中...")
        pngs = generate_thumbnails(pptx_path, thumb_dir)
        png_count = len(pngs)
        print(f"  サムネイル: {png_count}枚")
    else:
        print("  サムネイル: スキップ")

    # index.json 更新
    index = load_index()
    index["refs"] = [r for r in index["refs"] if r["id"] != ref_id]
    index["refs"].append({
        "id":            ref_id,
        "family":        family,
        "source_file":   str(pptx_path),
        "slide_count":   analysis["slide_count"],
        "analysis_path": str(analysis_path.relative_to(REFS_DIR.parent)),
        "thumbnail_dir": str(thumb_dir.relative_to(REFS_DIR.parent)),
        "analyzed_at":   analysis["analyzed_at"],
    })
    save_index(index)
    print(f"  index.json 更新完了")

    print(f"\n登録完了: {ref_id} ({analysis['slide_count']}枚, PNG:{png_count}枚)")
    return analysis


# ─── CLI ─────────────────────────────────────────────────

def slugify(text: str) -> str:
    text = re.sub(r"[^\w\s-]", "", text.lower())
    return re.sub(r"[\s_]+", "_", text)[:40].strip("_")


def cmd_extract_pattern(analysis_path: Path, slide_idx: int):
    """analysis.json から特定スライドのパターンを抽出"""
    if not analysis_path.exists():
        print(f"エラー: analysis.json が見つかりません: {analysis_path}")
        sys.exit(1)
    
    analysis = json.loads(analysis_path.read_text(encoding="utf-8"))
    slides = analysis.get("slides", [])
    
    if slide_idx >= len(slides):
        print(f"エラー: スライド {slide_idx} は存在しません（全 {len(slides)} 枚）")
        sys.exit(1)
    
    slide = slides[slide_idx]
    shapes = slide.get("shapes", [])
    
    print(f"\n{'='*80}")
    print(f"パターン抽出: {analysis_path.name} / スライド {slide_idx}")
    print(f"タイトル: {slide.get('title', '(なし)')}")
    print(f"{'='*80}\n")
    
    # 構造分析
    print(f"オブジェクト数: {len(shapes)}")
    print(f"レイアウト: {slide.get('layout_name', '不明')}\n")
    
    # 色の分布
    colors = [s.get("fill_color") for s in shapes if "fill_color" in s]
    if colors:
        from collections import Counter
        color_counts = Counter(colors)
        print("使用色:")
        for color, count in color_counts.most_common():
            print(f"  #{color}: {count}個")
        print()
    
    # オブジェクト一覧（座標・色・テキスト）
    print("オブジェクト詳細:\n")
    for i, shape in enumerate(shapes):
        print(f"[{i}] {shape.get('name', '(名前なし)')}")
        print(f"    type: {shape.get('shape_type', '不明')}")
        print(f"    left: {shape.get('left', 0):.3f}, top: {shape.get('top', 0):.3f}")
        print(f"    width: {shape.get('width', 0):.3f}, height: {shape.get('height', 0):.3f}")
        if "fill_color" in shape:
            print(f"    fill_color: #{shape['fill_color']}")
        if "text" in shape:
            print(f"    text: {shape['text']}")
        if "font_sizes" in shape:
            print(f"    font_sizes: {shape['font_sizes']}")
        print()


def main():
    parser = argparse.ArgumentParser(description="参照PPTXを分析してライブラリに登録")
    parser.add_argument("command", nargs="?", default="register", 
                        help="コマンド: register（デフォルト）, extract-pattern, list")
    parser.add_argument("pptx", nargs="?", help="分析するPPTXファイルパス or analysis.json")
    parser.add_argument("--id",           help="参照ID（省略時はファイル名から自動生成）")
    parser.add_argument("--family",       help="テンプレートファミリー（sx/jr等、省略時は親フォルダ名）")
    parser.add_argument("--no-thumbnail", action="store_true", help="サムネイル生成をスキップ")
    parser.add_argument("--slide",        type=int, help="抽出するスライド番号（extract-pattern用）")
    args = parser.parse_args()

    # コマンド判定（後方互換）
    command = args.command
    if command and Path(command).exists():
        # 引数1がファイルパス → 旧形式（register）
        args.pptx = command
        command = "register"

    if command == "list":
        index = load_index()
        if not index["refs"]:
            print("登録済みの参照作品はありません")
        else:
            for ref in index["refs"]:
                print(f"  [{ref['family']}] {ref['id']}  {ref['slide_count']}枚  {ref['analyzed_at']}")
                print(f"         {ref['source_file']}")
        return

    if command == "extract-pattern":
        if not args.pptx:
            print("エラー: analysis.json のパスを指定してください")
            parser.print_help()
            sys.exit(1)
        if args.slide is None:
            print("エラー: --slide <番号> を指定してください")
            sys.exit(1)
        cmd_extract_pattern(Path(args.pptx), args.slide)
        return

    # register (デフォルト)
    if not args.pptx:
        parser.print_help()
        sys.exit(1)

    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        print(f"エラー: ファイルが見つかりません: {pptx_path}")
        sys.exit(1)

    # family: 引数 → 親フォルダ名
    family = args.family or pptx_path.parent.name

    # id: 引数 → ファイル名から生成
    ref_id = args.id or f"{family}_{slugify(pptx_path.stem)}"

    register(pptx_path, ref_id, family, skip_thumbnail=args.no_thumbnail)


if __name__ == "__main__":
    main()
