"""
validate_slide.py
生成された PPTX のスライド内容を検証するスクリプト
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

def validate_pptx(pptx_path):
    """PPTX を検証して、オブジェクト配置・フォント・テキストを出力"""
    print(f"\n{'='*80}")
    print(f"PPTX 検証: {Path(pptx_path).name}")
    print(f"{'='*80}\n")
    
    prs = Presentation(pptx_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n--- スライド {slide_idx} ---")
        print(f"タイトル: {slide.shapes.title.text if slide.shapes.title else '(なし)'}\n")
        
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.name in ("Title", "Subtitle", "Body"):
                continue  # プレースホルダーはスキップ
            
            left = shape.left.inches
            top = shape.top.inches
            width = shape.width.inches
            height = shape.height.inches
            
            print(f"【Object {shape_idx}】")
            print(f"  Type: {shape.shape_type}")
            print(f"  Position: left={left:.2f}\", top={top:.2f}\"")
            print(f"  Size: width={width:.2f}\", height={height:.2f}\"")
            
            # テキストがある場合
            if hasattr(shape, "text_frame"):
                text = shape.text
                if text:
                    print(f"  Text: {text[:60]}{'...' if len(text) > 60 else ''}")
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        if para.runs:
                            run = para.runs[0]
                            font_size = run.font.size.pt if run.font.size else "default"
                            print(f"    [段 {para_idx}] font_size={font_size}pt")
            
            # 図形の塗りつぶし色
            if hasattr(shape, "fill") and shape.fill.type == 1:  # SOLID
                try:
                    rgb = shape.fill.fore_color.rgb
                    print(f"  Fill: #{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")
                except:
                    print(f"  Fill: (色取得失敗)")
            
            print()
    
    # テンプレート情報
    print(f"\nスライドサイズ: {prs.slide_width.inches:.2f}\" × {prs.slide_height.inches:.2f}\"")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("使用方法: python validate_slide.py <pptx_path>")
        sys.exit(1)
    
    validate_pptx(sys.argv[1])
