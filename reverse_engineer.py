"""
reverse_engineer.py
テンプレートPPTXを解析してレイアウト情報をレシピ化するツール

使い方:
  python reverse_engineer.py                          # 全レイアウト解析
  python reverse_engineer.py --layout 6               # 特定レイアウトのみ
  python reverse_engineer.py --save                   # recipes/template_layouts.json に保存
  python reverse_engineer.py --thumbnail              # 各レイアウトのサムネイル生成
"""

import json
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER

TEMPLATE_PATH = Path(__file__).parent / "SX_提案書_3.0_16x9.pptx"
RECIPES_DIR   = Path(__file__).parent / "recipes"
RECIPES_DIR.mkdir(exist_ok=True)

# プレースホルダータイプの名称マップ
PH_TYPE_NAME = {
    1:  "BODY",
    2:  "CENTER_TITLE",
    3:  "DATE",
    4:  "FOOTER",
    5:  "SLIDE_NUMBER",
    6:  "SUBTITLE",
    7:  "TABLE",
    8:  "CHART",
    9:  "ORG_CHART",
    10: "MEDIA_CLIP",
    11: "OBJECT",
    12: "PICTURE",
    13: "TITLE",
    14: "BITMAP",
    15: "VERTICAL_BODY",
    16: "VERTICAL_OBJECT",
    17: "VERTICAL_TITLE",
}


def inches(emu):
    return round(emu / 914400, 3)


def analyze_layout(layout, layout_idx: int) -> dict:
    """1レイアウトの詳細を解析"""
    result = {
        "layout_index": layout_idx,
        "name": layout.name,
        "placeholders": [],
        "other_shapes": [],
    }

    for shape in layout.shapes:
        left   = inches(shape.left   or 0)
        top    = inches(shape.top    or 0)
        width  = inches(shape.width  or 0)
        height = inches(shape.height or 0)

        if shape.is_placeholder:
            ph = shape.placeholder_format
            ph_type_val = ph.type.real if hasattr(ph.type, "real") else int(ph.type)
            ph_info = {
                "ph_idx":    ph.idx,
                "ph_type":   PH_TYPE_NAME.get(ph_type_val, str(ph.type)),
                "ph_type_id": ph_type_val,
                "name":      shape.name,
                "left":      left,
                "top":       top,
                "width":     width,
                "height":    height,
                "is_image":  ph_type_val == 12,  # PICTURE placeholder
            }
            # テキストの場合は既存テキストも取得
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    ph_info["default_text"] = text[:80]
            result["placeholders"].append(ph_info)
        else:
            shape_info = {
                "name":   shape.name,
                "type":   str(shape.shape_type),
                "left":   left,
                "top":    top,
                "width":  width,
                "height": height,
            }
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    shape_info["text"] = text[:80]
            result["other_shapes"].append(shape_info)

    # 推奨スキーマを生成
    result["recommended_schema"] = build_recommended_schema(result)
    return result


def build_recommended_schema(layout_info: dict) -> dict:
    """レイアウト情報から推奨JSONスキーマを生成"""
    ph_map = {ph["ph_idx"]: ph for ph in layout_info["placeholders"]}
    schema = {"type": "content", "notes": layout_info["name"]}

    # タイトル (ph_idx=0)
    if 0 in ph_map:
        schema["title"] = "スライドタイトル"

    # サブタイトル/キーメッセージ (ph_idx=13 が多い)
    for idx in [13, 6]:
        if idx in ph_map:
            schema["subtitle"] = "キーメッセージ"
            break

    # 本文 (ph_idx=10, 14, 1, 2)
    for idx in [10, 14, 1, 2]:
        if idx in ph_map:
            schema["body"] = "・箇条書き1\n・箇条書き2\n・箇条書き3"
            break

    # 画像プレースホルダー
    image_phs = [ph for ph in layout_info["placeholders"] if ph["is_image"]]
    if image_phs:
        schema["image_placeholders"] = [
            {"ph_idx": ph["ph_idx"], "left": ph["left"], "top": ph["top"],
             "width": ph["width"], "height": ph["height"]}
            for ph in image_phs
        ]
        schema["note_images"] = "このレイアウトには画像プレースホルダーあり。imagesフィールドで指定可能"

    return schema


def main():
    parser = argparse.ArgumentParser(description="PPTX Template Reverse Engineer")
    parser.add_argument("--layout", type=int, help="解析するレイアウトのインデックス（省略時は全件）")
    parser.add_argument("--save",   action="store_true", help="recipes/template_layouts.json に保存")
    parser.add_argument("--thumbnail", action="store_true", help="各レイアウトのサムネイルPPTXを生成")
    args = parser.parse_args()

    prs = Presentation(str(TEMPLATE_PATH))
    print(f"テンプレート: {TEMPLATE_PATH.name}")
    print(f"スライドサイズ: {inches(prs.slide_width)}\" x {inches(prs.slide_height)}\"")
    print(f"レイアウト数: {len(prs.slide_layouts)}\n")

    results = []
    target_layouts = [args.layout] if args.layout is not None else range(len(prs.slide_layouts))

    for idx in target_layouts:
        layout = prs.slide_layouts[idx]
        info = analyze_layout(layout, idx)
        results.append(info)

        print(f"─── Layout [{idx}] {info['name']} ───")
        print(f"  プレースホルダー数: {len(info['placeholders'])}")
        for ph in info['placeholders']:
            img_mark = " [IMAGE]" if ph["is_image"] else ""
            print(f"    ph_idx={ph['ph_idx']} ({ph['ph_type']}{img_mark})"
                  f"  pos=({ph['left']}\", {ph['top']}\")"
                  f"  size={ph['width']}\" x {ph['height']}\"")
        if info["other_shapes"]:
            print(f"  その他シェイプ数: {len(info['other_shapes'])}")
        schema = info["recommended_schema"]
        print(f"  推奨スキーマキー: {list(schema.keys())}")
        print()

    if args.save:
        out_path = RECIPES_DIR / "template_layouts.json"
        out_path.write_text(json.dumps(results, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"保存: {out_path}")

    if args.thumbnail:
        _generate_thumbnails(prs, results)

    return results


def _generate_thumbnails(prs, layout_infos):
    """各レイアウトにダミーテキストを入れたサムネイルPPTXを生成"""
    from pptx.util import Pt
    from copy import deepcopy
    thumb_dir = RECIPES_DIR / "thumbnails"
    thumb_dir.mkdir(exist_ok=True)

    thumb_prs = Presentation(str(TEMPLATE_PATH))
    # 既存スライドを削除
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    slide_id_list = thumb_prs.slides._sldIdLst
    for sld_id in list(slide_id_list):
        rId = sld_id.get(f'{{{r_ns}}}id')
        if rId:
            thumb_prs.part.rels.pop(rId)
        slide_id_list.remove(sld_id)

    for info in layout_infos:
        idx = info["layout_index"]
        layout = thumb_prs.slide_layouts[idx]
        slide = thumb_prs.slides.add_slide(layout)

        for ph_info in info["placeholders"]:
            ph_idx = ph_info["ph_idx"]
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.idx == ph_idx:
                    if shape.has_text_frame and not ph_info["is_image"]:
                        shape.text_frame.text = f"[{ph_info['ph_type']} idx={ph_idx}]"
                    break

    out_path = thumb_dir / "all_layouts_thumbnail.pptx"
    thumb_prs.save(str(out_path))
    print(f"サムネイルPPTX生成: {out_path}")
    import subprocess
    subprocess.Popen(["powershell", "-Command", f"Start-Process '{out_path}'"])


if __name__ == "__main__":
    main()
