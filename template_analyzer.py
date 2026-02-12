"""
template_analyzer.py
テンプレート解析ツール: pptx ファイルを解析し、profile.json のドラフトを生成する。

使い方:
  python template_analyzer.py analyze "path/to/template.pptx" --id "my_template" --name "My Template"
  python template_analyzer.py list
"""

import json
import shutil
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

TEMPLATES_DIR = Path(__file__).parent / "templates"


def emu_to_inches(emu_val) -> float:
    """EMU (English Metric Units) をインチに変換"""
    return round(emu_val / 914400, 3)


def analyze_template(pptx_path: Path) -> dict:
    """pptx ファイルを解析し、ドラフト profile データを返す"""
    prs = Presentation(str(pptx_path))

    slide_width = emu_to_inches(prs.slide_width)
    slide_height = emu_to_inches(prs.slide_height)

    layouts = []
    image_areas = {}

    for i, layout in enumerate(prs.slide_layouts):
        layout_info = {
            "index": i,
            "name": layout.name,
            "placeholders": [],
            "image_shapes": [],
        }

        for ph in layout.placeholders:
            layout_info["placeholders"].append({
                "ph_idx": ph.placeholder_format.idx,
                "name": ph.name,
                "left": emu_to_inches(ph.left),
                "top": emu_to_inches(ph.top),
                "width": emu_to_inches(ph.width),
                "height": emu_to_inches(ph.height),
            })

        for shape in layout.shapes:
            if not shape.is_placeholder and hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    coords = {
                        "left": emu_to_inches(shape.left),
                        "top": emu_to_inches(shape.top),
                        "width": emu_to_inches(shape.width),
                        "height": emu_to_inches(shape.height),
                        "text": text,
                    }
                    layout_info["image_shapes"].append(coords)
                    if "画像" in text or "image" in text.lower():
                        image_areas[str(i)] = {
                            "left": coords["left"],
                            "top": coords["top"],
                            "width": coords["width"],
                            "height": coords["height"],
                        }

        layouts.append(layout_info)

    # ヒューリスティックでレイアウト自動マッピング
    auto_layouts = _auto_detect_layouts(layouts)

    profile = {
        "id": "",
        "name": "",
        "file": "template.pptx",
        "slide_size": {"width": slide_width, "height": slide_height},
        "content_area": {
            "top": 1.5,
            "bottom": round(slide_height - 0.5, 1),
            "left": 0.5,
            "right": round(slide_width - 0.5, 1),
        },
        "layouts": auto_layouts,
        "image_areas": image_areas,
        "body_placeholder_search_order": [10, 14, 1, 2],
        "colors": {
            "primary": "",
            "secondary": "",
            "accent": "",
            "neutral": "",
            "background": "FFFFFF",
        },
        "color_semantics": {},
        "_raw_layouts": layouts,
    }

    return profile


def _auto_detect_layouts(layouts: list[dict]) -> dict:
    """レイアウト名からタイプを推測する"""
    result = {}
    keywords = {
        "title": ["タイトル", "ドキュメンテーション", "title"],
        "agenda": ["目次", "アジェンダ", "agenda", "index"],
        "chapter": ["チャプター", "chapter", "セクション", "section"],
        "content": ["コンテンツ", "content", "本文", "ヘッドライン"],
        "end": ["エンド", "end", "終了", "thank"],
    }

    for layout in layouts:
        name_lower = layout["name"].lower()
        for type_key, kw_list in keywords.items():
            for kw in kw_list:
                if kw.lower() in name_lower:
                    if type_key not in result:
                        phs = {
                            ph["name"]: ph["ph_idx"]
                            for ph in layout["placeholders"]
                        }
                        placeholder_map = {}
                        for ph in layout["placeholders"]:
                            idx = ph["ph_idx"]
                            if idx == 0:
                                placeholder_map["title"] = 0
                            elif idx == 1:
                                placeholder_map["subtitle"] = 1
                            elif idx in (10, 14):
                                placeholder_map["body"] = idx
                            elif idx == 13:
                                placeholder_map["subtitle"] = 13
                        result[type_key] = {
                            "index": layout["index"],
                            "placeholders": placeholder_map,
                            "_detected_name": layout["name"],
                        }
                    break

    return result


def register_template(pptx_path: Path, template_id: str, template_name: str) -> Path:
    """テンプレートを解析して templates/<id>/ に登録する"""
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"File not found: {pptx_path}")

    template_dir = TEMPLATES_DIR / template_id
    template_dir.mkdir(parents=True, exist_ok=True)

    # テンプレートファイルをコピー
    dest_pptx = template_dir / "template.pptx"
    shutil.copy2(str(pptx_path), str(dest_pptx))
    print(f"  Copied: {pptx_path.name} -> {dest_pptx}")

    # 解析
    profile = analyze_template(pptx_path)
    profile["id"] = template_id
    profile["name"] = template_name

    # profile.json 保存（_raw_layouts 付き、後でClaude Codeが確認・編集）
    profile_path = template_dir / "profile.json"
    profile_path.write_text(
        json.dumps(profile, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(f"  Created: {profile_path}")

    # design_guide.md スケルトン生成
    guide_path = template_dir / "design_guide.md"
    if not guide_path.exists():
        slide_w = profile["slide_size"]["width"]
        slide_h = profile["slide_size"]["height"]
        guide_content = f"""# {template_name} デザインガイド

> **使い方**: このテンプレートでスライドを設計する時に読む。
> 共通デザイン原則は `skills/design_principles.md` を参照。

---

## 1. スライドサイズと座標空間

```
幅: {slide_w}インチ  高さ: {slide_h}インチ
コンテンツエリア: top {profile['content_area']['top']}〜{profile['content_area']['bottom']}
左マージン: left >= {profile['content_area']['left']}
右端制限: left + width <= {profile['content_area']['right']}
```

---

## 2. レイアウト一覧

TODO: profile.json の layouts を確認し、各レイアウトの詳細を記述

---

## 3. カラーパレット

TODO: テンプレートのブランドカラーを調査して記述

---

## 4. 図解パターン集（具体的座標）

TODO: テンプレートに合った座標でパターンを記述
"""
        guide_path.write_text(guide_content, encoding="utf-8")
        print(f"  Created: {guide_path}")

    # 結果サマリー
    raw_layouts = profile.get("_raw_layouts", [])
    print(f"\n  Slide size: {profile['slide_size']['width']} x {profile['slide_size']['height']} inches")
    print(f"  Layouts found: {len(raw_layouts)}")
    auto_detected = {k: v.get('_detected_name', '') for k, v in profile['layouts'].items()}
    if auto_detected:
        print(f"  Auto-detected mappings:")
        for type_key, name in auto_detected.items():
            print(f"    {type_key} -> {name}")
    print(f"  Image areas found: {len(profile['image_areas'])}")
    print(f"\n  Next steps:")
    print(f"    1. Review {profile_path}")
    print(f"    2. Confirm layout mappings (remove _detected_name, adjust placeholders)")
    print(f"    3. Remove _raw_layouts from profile.json")
    print(f"    4. Fill in colors and design_guide.md")

    return template_dir


def list_templates():
    """登録済みテンプレートを一覧表示"""
    if not TEMPLATES_DIR.exists():
        print("No templates directory found.")
        return

    templates = sorted(TEMPLATES_DIR.iterdir())
    if not templates:
        print("No templates registered.")
        return

    print(f"Registered templates ({TEMPLATES_DIR}):\n")
    for d in templates:
        if d.is_dir():
            profile_path = d / "profile.json"
            if profile_path.exists():
                data = json.loads(profile_path.read_text(encoding="utf-8"))
                name = data.get("name", "?")
                size = data.get("slide_size", {})
                print(f"  {d.name}: {name} ({size.get('width', '?')}x{size.get('height', '?')})")
            else:
                print(f"  {d.name}: (no profile.json)")


def main():
    parser = argparse.ArgumentParser(description="Template Analyzer for PPTX Generation")
    sub = parser.add_subparsers(dest="command")

    # analyze / register
    analyze_parser = sub.add_parser("analyze", help="Analyze a PPTX template and register it")
    analyze_parser.add_argument("pptx_path", help="Path to the PPTX template file")
    analyze_parser.add_argument("--id", required=True, help="Template ID (e.g., 'client_acme')")
    analyze_parser.add_argument("--name", required=True, help="Template display name")

    # list
    sub.add_parser("list", help="List registered templates")

    args = parser.parse_args()

    if args.command == "analyze":
        print(f"Analyzing template: {args.pptx_path}")
        register_template(Path(args.pptx_path), args.id, args.name)
    elif args.command == "list":
        list_templates()
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
