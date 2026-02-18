#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SX テンプレートの content layout の配置要件を分析
"""

from pptx import Presentation
from pathlib import Path
import json

template_path = Path("templates/sx_proposal/template.pptx")
prs = Presentation(str(template_path))

# content layout (index 6)
layout = prs.slide_layouts[6]

print("\n" + "="*70)
print("SX Template - Content Layout (index 6) Analysis")
print("="*70)

print(f"\nSlide Size: {prs.slide_width.inches:.3f}\" × {prs.slide_height.inches:.3f}\"")
print(f"\nPlaceholders in content layout:")

layout_config = {
    "slide_size": {
        "width": prs.slide_width.inches,
        "height": prs.slide_height.inches
    },
    "placeholders": {}
}

for placeholder in layout.placeholders:
    name = placeholder.name
    idx = placeholder.placeholder_format.idx
    left = placeholder.left.inches
    top = placeholder.top.inches
    width = placeholder.width.inches
    height = placeholder.height.inches
    
    print(f"\n  [{idx}] {name}")
    print(f"      Position: ({left:.3f}\", {top:.3f}\")")
    print(f"      Size: {width:.3f}\" × {height:.3f}\"")
    print(f"      Bottom-Right: ({left + width:.3f}\", {top + height:.3f}\")")
    
    layout_config["placeholders"][name] = {
        "index": idx,
        "left": round(left, 3),
        "top": round(top, 3),
        "width": round(width, 3),
        "height": round(height, 3),
        "right": round(left + width, 3),
        "bottom": round(top + height, 3)
    }

# Content area (from profile.json)
print("\n" + "-"*70)
print("Content Area (from profile.json):")
print(f"  left: 0.5\", top: 1.5\", right: 12.8\", bottom: 7.0\"")
print(f"  Available width: 12.3\", Available height: 5.5\"")

print("\n" + "-"*70)
print("Canvas Designer Configuration:")
print(f"  Content area within which objects should be placed:")
print(f"    - Horizontal: 0.5\" to 12.8\" (12.3\" wide)")
print(f"    - Vertical: 1.5\" to 7.0\" (5.5\" high)")

# 保存
with open("sx_template_layout_analysis.json", "w", encoding="utf-8") as f:
    json.dump(layout_config, f, indent=2, ensure_ascii=False)

print(f"\n✅ Analysis saved to: sx_template_layout_analysis.json")
print("="*70)
