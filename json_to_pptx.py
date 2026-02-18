#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
デモスライドを PPTX に変換
汎用ツールで生成した JSON から PowerPoint を生成
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def hex_to_rgb(hex_color):
    """16進数カラーコード を RGB に変換"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def add_box(slide, obj):
    """背景付きテキストボックスを追加"""
    left = Inches(obj['left'])
    top = Inches(obj['top'])
    width = Inches(obj['width'])
    height = Inches(obj['height'])
    
    shape = slide.shapes.add_shape(
        1,  # rectangle
        left, top, width, height
    )
    
    # 背景色
    fill_color = obj.get('fill_color', 'FFFFFF')
    rgb = hex_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)
    
    # テキスト
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = 1  # middle
    
    p = text_frame.paragraphs[0]
    p.text = obj.get('text', '')
    p.alignment = PP_ALIGN.CENTER
    
    # テキスト色
    font_color = obj.get('font_color', '000000')
    rgb_font = hex_to_rgb(font_color)
    p.font.color.rgb = RGBColor(*rgb_font)
    
    # フォントサイズ
    p.font.size = Pt(obj.get('font_size', 12))
    p.font.name = 'Arial'
    
    # 枠線なし
    shape.line.color.rgb = RGBColor(*rgb)

def add_text(slide, obj):
    """テキストボックスを追加（背景なし）"""
    left = Inches(obj['left'])
    top = Inches(obj['top'])
    width = Inches(obj['width'])
    height = Inches(obj['height'])
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    # 複数行対応
    text = obj.get('text', '')
    lines = text.split('\n')
    
    for i, line in enumerate(lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = line
        p.level = 0
        
        # テキスト色
        font_color = obj.get('font_color', '000000')
        rgb_font = hex_to_rgb(font_color)
        p.font.color.rgb = RGBColor(*rgb_font)
        
        # フォントサイズ
        p.font.size = Pt(obj.get('font_size', 11))
        p.font.name = 'Arial'

def json_to_pptx(json_path, output_path):
    """JSON ファイルから PPTX を生成"""
    
    # JSON を読み込み
    with open(json_path, 'r', encoding='utf-8') as f:
        slide_data = json.load(f)
    
    # プレゼンテーションを作成
    prs = Presentation()
    prs.slide_width = Inches(12.8)
    prs.slide_height = Inches(7.2)
    
    # スライドを追加（白背景）
    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # オブジェクトを追加
    for obj in slide_data.get('objects', []):
        obj_type = obj.get('type')
        
        if obj_type == 'box':
            add_box(slide, obj)
        elif obj_type == 'text':
            add_text(slide, obj)
        elif obj_type == 'arrow':
            # 矢印は簡易的に線で表現
            left = Inches(obj['left'])
            top = Inches(obj['top'])
            width = Inches(obj['width'])
            height = Inches(obj['height'])
            
            connector = slide.shapes.add_connector(1, left, top, left + width, top)
            line = connector.line
            line.color.rgb = RGBColor(*hex_to_rgb(obj.get('fill_color', '404040')))
            line.width = Pt(3)
    
    # PPTX を保存
    prs.save(output_path)
    return output_path

def main():
    print("\n" + "="*70)
    print("JSON → PPTX 変換")
    print("="*70)
    
    json_files = [
        ("test_output/demo1_comparison.json", "output_demo1_comparison.pptx"),
        ("test_output/demo2_three_tier.json", "output_demo2_three_tier.pptx"),
        ("test_output/demo3_custom.json", "output_demo3_custom.pptx"),
    ]
    
    for json_path, pptx_path in json_files:
        json_file = Path(json_path)
        if not json_file.exists():
            print(f"\n❌ {json_path} が見つかりません")
            continue
        
        print(f"\n[*] {json_path} を処理中...")
        
        try:
            with open(json_path, 'r', encoding='utf-8') as f:
                slide_data = json.load(f)
            
            title = slide_data.get('title', 'Untitled')
            print(f"    タイトル: {title}")
            print(f"    オブジェクト数: {len(slide_data.get('objects', []))}")
            
            output = json_to_pptx(json_path, pptx_path)
            print(f"    ✅ PPTX生成: {pptx_path}")
        
        except Exception as e:
            print(f"    ❌ エラー: {e}")
    
    print("\n" + "="*70)
    print("✅ PPTX 生成完了")
    print("="*70)
    print(f"\n📋 生成されたファイル:")
    print(f"   • output_demo1_comparison.pptx")
    print(f"   • output_demo2_three_tier.pptx")
    print(f"   • output_demo3_custom.pptx")
    print(f"\n🔍 確認方法:")
    print(f"   Windows: ファイルエクスプローラーで .pptx を開く")
    print(f"   または PowerPoint で開く")
    print(f"\n" + "="*70)

if __name__ == "__main__":
    main()
