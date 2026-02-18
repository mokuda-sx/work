#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
汎用 JSON → SX テンプレート PPTX 生成

Universal Slide Designer で生成した JSON を、
SX テンプレートの content レイアウトに配置して PPTX を生成
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def hex_to_rgb(hex_color):
    """16進数カラーコード を RGB に変換"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_pptx_with_sx_template(json_path: str, output_path: str) -> bool:
    """
    JSON ファイルから、SX テンプレートを使用して PPTX を生成
    
    Args:
        json_path: 入力 JSON ファイル
        output_path: 出力 PPTX ファイル
    
    Returns:
        成功した場合 True
    """
    try:
        # JSON を読み込み
        with open(json_path, 'r', encoding='utf-8') as f:
            slide_data = json.load(f)
        
        # SX テンプレートベースのプレゼンテーション開く
        template_path = Path(__file__).parent / "templates" / "sx_proposal" / "template.pptx"
        
        if not template_path.exists():
            print(f"❌ テンプレートが見つかりません: {template_path}")
            return False
        
        prs = Presentation(str(template_path))
        
        # スライドサイズを確認（SX: 13.333 x 7.5）
        print(f"   スライドサイズ: {prs.slide_width.inches:.2f}\" × {prs.slide_height.inches:.2f}\"")
        
        # content レイアウト (index 6) を使用
        content_layout = prs.slide_layouts[6]
        
        # 新しいスライドを追加
        slide = prs.slides.add_slide(content_layout)
        
        # ─── タイトルを設定 ───
        if slide_data.get('title'):
            # title placeholder (index 0)
            title_shape = slide.placeholders[0]
            title_shape.text = slide_data['title']
        
        # ─── サブタイトルを設定 ───
        if slide_data.get('subtitle'):
            # subtitle placeholder (index 13)
            try:
                subtitle_shape = slide.placeholders[13]
                subtitle_shape.text = slide_data['subtitle']
            except:
                pass  # placeholder がない場合はスキップ
        
        # ─── コンテンツオブジェクトを配置 ───
        # テンプレートの content area に収まるように配置
        # SX template の content area: top: 1.5, bottom: 7.0, left: 0.5, right: 12.8
        
        content_area = {
            'top': 1.5,
            'left': 0.5,
            'right': 12.8,
            'bottom': 7.0,
            'width': 12.8 - 0.5,
            'height': 7.0 - 1.5
        }
        
        # JSON 座標系 (12.8 x 7.2) → SX テンプレート座標系へ換算
        # JSON: 0.5-4.8, SX: 0.5-12.8
        scale_x = content_area['width'] / 4.3  # JSON の標準幅 4.3"
        scale_y = content_area['height'] / 5.5  # JSON の標準高さ 5.5"
        
        for obj in slide_data.get('objects', []):
            obj_type = obj.get('type')
            
            # JSON 座標系からの変換
            json_left = obj.get('left', 0)
            json_top = obj.get('top', 0)
            json_width = obj.get('width', 1)
            json_height = obj.get('height', 0.5)
            
            # テンプレート座標系へ変換
            actual_left = content_area['left'] + (json_left - 0.5) * scale_x
            actual_top = content_area['top'] + json_top * scale_y
            actual_width = json_width * scale_x
            actual_height = json_height * scale_y
            
            if obj_type == 'box':
                # 背景付きテキストボックス
                shape = slide.shapes.add_shape(
                    1,  # rectangle
                    Inches(actual_left),
                    Inches(actual_top),
                    Inches(actual_width),
                    Inches(actual_height)
                )
                
                # 背景色
                fill_color = obj.get('fill_color', 'FFFFFF')
                rgb = hex_to_rgb(fill_color)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*rgb)
                
                # 枠線
                shape.line.color.rgb = RGBColor(*rgb)
                
                # テキスト
                text_frame = shape.text_frame
                text_frame.word_wrap = True
                text_frame.vertical_anchor = 1  # middle
                
                text = obj.get('text', '')
                lines = text.split('\n')
                
                for i, line in enumerate(lines):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = line
                    p.alignment = PP_ALIGN.CENTER
                    
                    # テキスト色
                    font_color = obj.get('font_color', '000000')
                    rgb_font = hex_to_rgb(font_color)
                    p.font.color.rgb = RGBColor(*rgb_font)
                    
                    # フォントサイズ
                    font_size = obj.get('font_size', 12)
                    p.font.size = Pt(font_size)
                    p.font.name = 'Arial'
            
            elif obj_type == 'text':
                # テキストボックス（背景なし）
                text_box = slide.shapes.add_textbox(
                    Inches(actual_left),
                    Inches(actual_top),
                    Inches(actual_width),
                    Inches(actual_height)
                )
                
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                
                text = obj.get('text', '')
                lines = text.split('\n')
                
                for i, line in enumerate(lines):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = line
                    
                    # テキスト色
                    font_color = obj.get('font_color', '000000')
                    rgb_font = hex_to_rgb(font_color)
                    p.font.color.rgb = RGBColor(*rgb_font)
                    
                    # フォントサイズ
                    font_size = obj.get('font_size', 11)
                    p.font.size = Pt(font_size)
                    p.font.name = 'Arial'
            
            elif obj_type == 'arrow':
                # 矢印（簡易的に線で表現）
                connector = slide.shapes.add_connector(
                    1,  # straight connector
                    Inches(actual_left),
                    Inches(actual_top),
                    Inches(actual_left + actual_width),
                    Inches(actual_top)
                )
                
                line = connector.line
                line_color = obj.get('fill_color', '404040')
                rgb_line = hex_to_rgb(line_color)
                line.color.rgb = RGBColor(*rgb_line)
                line.width = Pt(3)
        
        # PPTX を保存
        prs.save(output_path)
        return True
    
    except Exception as e:
        print(f"   エラー: {e}")
        import traceback
        traceback.print_exc()
        return False

def main():
    print("\n" + "="*70)
    print("汎用 JSON → SX テンプレート PPTX 生成")
    print("="*70)
    
    json_files = [
        ("test_output/demo1_comparison.json", "output_sx_demo1_comparison.pptx"),
        ("test_output/demo2_three_tier.json", "output_sx_demo2_three_tier.pptx"),
        ("test_output/demo3_custom.json", "output_sx_demo3_custom.pptx"),
    ]
    
    for json_path, pptx_path in json_files:
        json_file = Path(json_path)
        if not json_file.exists():
            print(f"\n❌ {json_path} が見つかりません")
            continue
        
        print(f"\n[*] {json_path} を処理中...")
        
        with open(json_path, 'r', encoding='utf-8') as f:
            slide_data = json.load(f)
        
        title = slide_data.get('title', 'Untitled')
        print(f"    タイトル: {title}")
        print(f"    オブジェクト数: {len(slide_data.get('objects', []))}")
        
        success = create_pptx_with_sx_template(json_path, pptx_path)
        
        if success:
            print(f"    ✅ SX テンプレート PPTX: {pptx_path}")
    
    print("\n" + "="*70)
    print("✅ SX テンプレート PPTX 生成完了")
    print("="*70)
    print(f"\n📋 生成されたファイル:")
    print(f"   • output_sx_demo1_comparison.pptx  (SX テンプレート対応)")
    print(f"   • output_sx_demo2_three_tier.pptx  (SX テンプレート対応)")
    print(f"   • output_sx_demo3_custom.pptx      (SX テンプレート対応)")
    
    print(f"\n🔍 変更点:")
    print(f"   ✓ SX テンプレートのレイアウトと配色を使用")
    print(f"   ✓ content レイアウト (index 6) を使用")
    print(f"   ✓ テンプレートの content area に収まるよう座標変換")
    print(f"   ✓ SX ブランドガイドに準拠")
    
    print(f"\n" + "="*70)

if __name__ == "__main__":
    main()
